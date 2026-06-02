"""
helpers.py — Galledia Präsentations-Skill v1.7 (CI-Grössen 72/30pt + Längenlimit 32/35 Zeichen)
Schrift-Regel: Volte (Fliesstext/Kapiteltitel), Volte Semibold (Headlines/Hervorhebungen)
Keine Versalien. Kein Volte Rounded in generiertem Inhalt.
"""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── CI-Farb-Tokens ────────────────────────────────────────────────────────────
RED  = RGBColor(0xE6,0x1C,0x52); BLACK = RGBColor(0,0,0); WHITE = RGBColor(0xFF,0xFF,0xFF)
G1   = RGBColor(0x40,0x40,0x40); G2    = RGBColor(0x66,0x66,0x66)
G3   = RGBColor(0xA6,0xA6,0xA6); G4    = RGBColor(0xD9,0xD9,0xD9)
GL   = RGBColor(0xF2,0xF2,0xF5); TURK  = RGBColor(0x22,0xAA,0x9F)

# ── CI-Schrift-Tokens ─────────────────────────────────────────────────────────
BODY = "Volte"           # Fliesstext, Kapiteltitel, Labels
SB   = "Volte Semibold"  # Headlines, Kartenüberschriften, Hervorhebungen

# Titel-Grössen = CI-Vorgabe (fix, nie adaptiv). Länge begrenzt → immer 1 Zeile.
HEAD_PT  = 72   # Headline — CI-Vorgabe
KICK_PT  = 30   # Kapiteltitel — CI-Vorgabe
MAX_HEAD = 32   # max. Zeichen Headline bei 72pt (gemessen: ~32 passen auf 1 Zeile)
MAX_KICK = 35   # max. Zeichen Kapiteltitel bei 30pt

# ── Template-Grid (aus Vorlage_6.pptx, idx=0/11 auf 04_vielText) ──────────────
GX   = Inches(0.979)   # Linke Hilfslinie
GY_K = Inches(0.750)   # Y Kapiteltitel
GY_H = Inches(1.061)   # Y Headline
GW   = Inches(18.026)  # Content-Breite
SH   = Inches(11.25)   # Folienhöhe

LAYOUTS = {
    "titel":"Titelfolie","abschluss":"Abschlussfolie","schluss":"Schlussfolie",
    "zwischenrot":"2_Zwischenfolie rot","agenda5":"01_Agenda 5","agenda22":"01_Agenda 22",
    "wenig":"02_wenigText","viel":"04_vielText","leer":"Leer","piktogramme":"Piktogramme",
}
_TEMPLATE = os.path.join(os.path.dirname(__file__),"assets","Vorlage_6.pptx")
_LOGO_DIR  = os.path.join(os.path.dirname(__file__),"assets","logo")
LOGOS = {  # CI: Rot→helle Folien | Weiss→rote/dunkle Folien | Schwarz→fehlt (neu exportieren)
    "rot":              os.path.join(_LOGO_DIR,"logo_rot.png"),
    "rot_schriftzug":   os.path.join(_LOGO_DIR,"logo_rot_schriftzug.png"),
    "weiss":            os.path.join(_LOGO_DIR,"logo_weiss.png"),
    "weiss_schriftzug":   os.path.join(_LOGO_DIR,"logo_weiss_schriftzug.png"),
    "schwarz":            os.path.join(_LOGO_DIR,"logo_schwarz.png"),
    "schwarz_schriftzug": os.path.join(_LOGO_DIR,"logo_schwarz_schriftzug.png"),
}

# ── Basis-Setup ───────────────────────────────────────────────────────────────

def build_presentation(template=None, datum="", rechtseinheit=""):
    """
    Vorlage laden, Beispielfolien entfernen.
    datum / rechtseinheit werden als Fusszeilen-Defaults gespeichert.
    Beispiel: build_presentation(datum='29. Mai 2026', rechtseinheit='Galledia Fachmedien AG')
    """
    prs = Presentation(template or _TEMPLATE)
    sl = prs.slides._sldIdLst
    for sid in list(sl):
        prs.part.drop_rel(sid.get(qn("r:id"))); sl.remove(sid)
    prs._layouts = {l.name:l for l in prs.slide_masters[0].slide_layouts}
    prs._datum = datum
    prs._rechtseinheit = rechtseinheit
    _strip_transitions(prs)
    return prs


def _strip_transitions(prs):
    """Entfernt alle Folienübergänge aus Master, Layouts und Slides."""
    TAG = qn('p:transition')
    sources = []
    for m in prs.slide_masters:
        sources.append(m._element)
        for lay in m.slide_layouts:
            sources.append(lay._element)
    for slide in prs.slides:
        sources.append(slide._element)
    for el in sources:
        for tr in el.findall('.//' + TAG):
            tr.getparent().remove(tr)


def _move_placeholder(ph, x_emu, y_emu, w_emu, h_emu):
    """Setzt Position+Grösse eines Platzhalters explizit (überschreibt Layout-Vererbung)."""
    from lxml import etree
    sp = ph._element
    sp_pr = sp.find(qn('p:spPr'))
    if sp_pr is None:
        sp_pr = etree.SubElement(sp, qn('p:spPr'))
    xfrm = sp_pr.find(qn('a:xfrm'))
    if xfrm is None:
        xfrm = etree.SubElement(sp_pr, qn('a:xfrm'))
    off = xfrm.find(qn('a:off'))
    if off is None:
        off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(int(x_emu))); off.set('y', str(int(y_emu)))
    ext = xfrm.find(qn('a:ext'))
    if ext is None:
        ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(int(w_emu))); ext.set('cy', str(int(h_emu)))


def _lay(prs, key):
    name = LAYOUTS.get(key, key)
    cache = getattr(prs,"_layouts",{}) or {l.name:l for l in prs.slide_masters[0].slide_layouts}
    if name not in cache:
        raise ValueError(f"Layout '{name}' fehlt. Verfügbar: {list(cache)}")
    return cache[name]

def _blank(prs): return prs.slides.add_slide(_lay(prs,"leer"))

def _phs(slide): return {p.placeholder_format.idx:p for p in slide.placeholders}

# ── Primitive ─────────────────────────────────────────────────────────────────

def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    b=slide.shapes.add_textbox(x,y,w,h); tf=b.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    return tf

def run(para, text, size, font, color,
        align=PP_ALIGN.LEFT, space_after=6, bold=False, italic=False):
    # Defensive: leerer Text erzeugt keinen run → IndexError. Mindestens 1 Zeichen sichern.
    para.text = text if text else " "
    para.alignment = align; para.space_after = Pt(space_after)
    if not para.runs:  # extra safety net
        para.add_run().text = " "
    r = para.runs[0]; r.font.name = font; r.font.size = Pt(size)
    r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic
    return para

_MD_BOLD = re.compile(r'\*\*(.+?)\*\*')

def _strip_md(s):
    """Entfernt **bold**-Markdown — Bold rendern wir über Schrift, nicht über Marker."""
    return _MD_BOLD.sub(r'\1', str(s or ''))

def _check_len(text, limit, was):
    """Erzwingt CI-Längenlimit. Zu lang → klarer Fehler (zwingt zu knackigem Titel)."""
    t = _strip_md(text).strip()
    if len(t) > limit:
        raise ValueError(
            f"{was} zu lang: {len(t)} Zeichen (max. {limit}). "
            f"Bitte kürzer formulieren: «{t}»")
    return t

def _headline_size(text):
    """Adaptive Headline-Grösse — verhindert Überlauf in den Body bei langen Headlines."""
    n = len(text or "")
    if n <= 28: return 72
    if n <= 38: return 54
    if n <= 50: return 44
    if n <= 64: return 36
    return 30

def _enable_shrink(tf):
    """Aktiviert 'Text bei Überlauf verkleinern' (normAutofit) auf einem Textframe."""
    from lxml import etree
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        bodyPr = etree.SubElement(tf._txBody, qn('a:bodyPr'))
    for tag in ('a:spAutoFit','a:noAutofit','a:normAutofit'):
        e = bodyPr.find(qn(tag))
        if e is not None: bodyPr.remove(e)
    etree.SubElement(bodyPr, qn('a:normAutofit'))

def card(slide, x, y, w, h, fill, corner=0.06):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.fill.background(); sh.shadow.inherit=False
    try: sh.adjustments[0]=corner
    except: pass
    return sh

# ── Interne Header/Footer-Helfer ──────────────────────────────────────────────

def _footer_str(prs, folio):
    """Baut den Fusszeilen-Text: 'Folie [Nr], [Datum], [Rechtseinheit]'"""
    parts = []
    if folio: parts.append(f"Folie {folio}")
    if getattr(prs,"_datum",""): parts.append(prs._datum)
    if getattr(prs,"_rechtseinheit",""): parts.append(prs._rechtseinheit)
    return ", ".join(parts)

def _set_header(slide, kicker, headline):
    """
    Kapiteltitel + Headline — EINHEITLICH über alle Folientypen.
    Native Titel-Platzhalter werden geleert; Text kommt in manuelle Textboxen
    an exakt denselben Koordinaten, Grössen und Anker. Garantiert identische
    Position und Schriftgrösse auf jeder Folie. Headline-Box ist hoch genug
    für 2 Zeilen — kein Autofit, keine adaptive Skalierung.
    """
    ph = _phs(slide)
    kicker = _check_len(kicker, MAX_KICK, "Kapiteltitel")
    headline = _check_len(headline, MAX_HEAD, "Headline")
    for idx in (0, 11):
        if idx in ph:
            try: ph[idx].text = ""
            except: pass
    # Kapiteltitel — 30pt CI
    tfk = tb(slide, GX, GY_K, GW, Inches(0.40), MSO_ANCHOR.TOP)
    run(tfk.paragraphs[0], kicker, KICK_PT, BODY, BLACK, space_after=0)
    # Headline — 72pt CI, 1 Zeile (durch Längenlimit garantiert)
    if headline:
        tfh = tb(slide, GX, GY_H, GW, Inches(1.20), MSO_ANCHOR.TOP)
        run(tfh.paragraphs[0], headline, HEAD_PT, SB, BLACK, space_after=0)
    if 13 in ph and not ph[13].text_frame.text:
        ph[13].text=""   # Prompt-Text unterdrücken

def _set_footer(slide, prs, folio, source=""):
    """Fusszeile links: 'Folie [Nr], [Datum], [Rechtseinheit]'. Quellenangabe darüber."""
    ph = _phs(slide)
    footer = _footer_str(prs, folio)
    if footer:
        if 14 in ph:
            run(ph[14].text_frame.paragraphs[0], footer, 11, BODY, BLACK,
                align=PP_ALIGN.LEFT, space_after=0)
        else:
            # Textbox auf exakter Hilfslinie (y=10.82", h=0.25", w=GW — wie idx=14 auf 04_vielText)
            tf=tb(slide,GX,Inches(10.82),GW,Inches(0.25))
            run(tf.paragraphs[0], footer, 11, BODY, BLACK, space_after=0)
    if source:
        tf=tb(slide,GX,SH-Inches(1.2),GW,Inches(0.35))
        run(tf.paragraphs[0],f"Quelle: {source}",11,BODY,G2,italic=True,space_after=0)

# ── Native-Layout-Slides ──────────────────────────────────────────────────────

# Titelfolie-Geometrie (aus Vorlage_6 gemessen)
_TITLE_Y_SUB        = Inches(3.968)   # y idx=10 Untertitel (bei 2-zeiligem Haupttitel)
_TITLE_Y_MAIN       = Inches(4.683)   # y idx=0  Haupttitel
_TITLE_1LINE_SHIFT  = Inches(0.68)    # Verschiebung nach unten bei 1-zeiligem Haupttitel
_TITLE_CHARS_1LINE  = 40              # Schwellenwert: ≤ 1 Zeile; > 2 Zeilen

def add_title(prs, title, subtitle=""):
    """
    Titelfolie. Dynamische Zentrierung:
    Einzeiliger Haupttitel (≤40 Zeichen) → beide Platzhalter um 0.68" nach unten,
    sodass der Block optisch zentriert bleibt.
    """
    s=prs.slides.add_slide(_lay(prs,"titel"))
    ph=_phs(s); ph[0].text=title
    if 10 in ph and subtitle: ph[10].text=subtitle
    if len(title) <= _TITLE_CHARS_1LINE:
        shift = int(_TITLE_1LINE_SHIFT)
        # Gemessene Layout-Koordinaten (aus Vorlage_6): x, y, w, h in EMU
        specs = {
            0:  (Inches(1.023), Inches(4.683), Inches(17.955), Inches(2.175)),
            10: (Inches(1.023), Inches(3.968), Inches(17.955), Inches(0.752)),
        }
        for idx in (0, 10):
            if idx in ph:
                x, y, w, h = specs[idx]
                _move_placeholder(ph[idx], x, int(y) + shift, w, h)
    return s

def add_section(prs, number, title):
    """Optionaler Kapitel-Anker (für Decks >15 Folien mit mehreren Kapiteln)."""
    s=prs.slides.add_slide(_lay(prs,"zwischenrot"))
    tf=tb(s,GX,Inches(3.8),GW,Inches(4),MSO_ANCHOR.TOP)
    run(tf.paragraphs[0], number, 54, SB, WHITE, space_after=4)
    run(tf.add_paragraph(), title,  40, SB, WHITE)
    return s

def add_agenda(prs, items, variant="agenda5", folio=""):
    """
    Agenda-Folie. «Agenda» steht in Kapiteltitel (idx=0) UND Haupttitel (idx=11).
    variant: 'agenda5' (≤5 Punkte) oder 'agenda22' (6–12 Punkte).
    """
    s=prs.slides.add_slide(_lay(prs,variant)); ph=_phs(s)
    ph[0].text="Agenda"; ph[11].text="Agenda"
    tf=ph[13].text_frame; tf.paragraphs[0].text=items[0]
    for item in items[1:]: tf.add_paragraph().text=item
    if 14 in ph: run(ph[14].text_frame.paragraphs[0],
                     _footer_str(prs,folio),11,BODY,BLACK,space_after=0)
    return s

def _set_bullet(para, char, font="Arial", indent_in=0.0):
    """Setzt Bullet-Zeichen + hängenden Einzug auf einen Absatz."""
    pPr = para._p.get_or_add_pPr()
    if char is None:
        pPr.set('marL', '0'); pPr.set('indent', '0')
    else:
        marL = int(Inches(indent_in + 0.28))
        pPr.set('marL', str(marL))
        pPr.set('indent', str(-int(Inches(0.28))))
    for tag in ('a:buNone','a:buChar','a:buAutoNum','a:buFont'):
        e = pPr.find(qn(tag))
        if e is not None: pPr.remove(e)
    if char is None:
        pPr.append(pPr.makeelement(qn('a:buNone'), {}))
        return
    pPr.append(pPr.makeelement(qn('a:buFont'), {'typeface': font}))
    pPr.append(pPr.makeelement(qn('a:buChar'), {'char': char}))

# Bullet-Marker (alles wird akzeptiert; rendert immer als CI-Bullet •/–):
#   Ebene 1:  "· ", "• ", "- ", "* "        (auch "•  ", "*  " mit Doppel-Space)
#   Ebene 2:  "·· ", "•• ", "-- ", "  - ", "  * "
_BUL_L2 = re.compile(r'^(?:··|••|--|  [-*•·])\s+')
_BUL_L1 = re.compile(r'^[·•\-*]\s+')

def _render_body(tf, body_text):
    """
    Rendert body_text mit echter Formatierung statt rohem Text-Dump:
      Bullet-Marker am Zeilenanfang  → Ebene-1-Bullet (•, schwarz, Volte 19pt Regular)
      Doppel-Marker / Einrückung      → Ebene-2-Bullet (–, grau, Volte 17pt, eingerückt)
      Zeile mit `**text**` oder ohne Marker → Zwischentitel (Volte Semibold 22pt, schwarz)
      Leerzeile                       → kleiner Abstand
    Markdown `**bold**` wird überall gestrippt — Bold rendern wir via Font, nicht Marker.
    """
    lines = [l for l in str(body_text).split("\n")]
    tf.word_wrap = True
    first = True
    def _para():
        nonlocal first
        if first:
            first = False
            return tf.paragraphs[0]
        return tf.add_paragraph()
    for raw in lines:
        s = raw.rstrip()
        # Leerzeile
        if not s.strip():
            p = _para(); p.space_after = Pt(4); _set_bullet(p, None); continue
        # Ebene-2-Bullet zuerst prüfen (Längere Marker)
        m2 = _BUL_L2.match(s)
        if m2:
            text = _strip_md(s[m2.end():]).strip()
            p = _para(); run(p, text, 17, BODY, G2, space_after=5)
            p.level = 1; _set_bullet(p, "–", indent_in=0.4); continue
        # Ebene-1-Bullet
        m1 = _BUL_L1.match(s.lstrip())
        if m1 and not s.startswith(" "):  # nicht eingerückt = L1
            text = _strip_md(s.lstrip()[m1.end():]).strip()
            p = _para(); run(p, text, 19, BODY, BLACK, space_after=6)
            p.level = 0; _set_bullet(p, "•"); continue
        # Zwischentitel: Zeile ohne Bullet-Marker (egal ob mit ** oder ohne)
        text = _strip_md(s).strip()
        p = _para(); run(p, text, 22, SB, BLACK, space_after=4)
        p.space_before = Pt(10); _set_bullet(p, None)
    _enable_shrink(tf)

def add_content(prs, variant, kapitel, headline, body_text, folio="", source=""):
    """Inhaltsfolie. variant='viel' oder 'wenig'."""
    s=prs.slides.add_slide(_lay(prs,variant)); ph=_phs(s)
    _set_header(s, kapitel, headline)
    _render_body(ph[13].text_frame, body_text)
    if 14 in ph: run(ph[14].text_frame.paragraphs[0],
                     _footer_str(prs,folio),11,BODY,BLACK,space_after=0)
    if source and 15 in ph: ph[15].text=f"Quelle: {source}"
    return s

def add_closing(prs):   return prs.slides.add_slide(_lay(prs,"schluss"))
def add_discussion(prs): return prs.slides.add_slide(_lay(prs,"abschluss"))

def add_logo(slide, variant="rot", size_in=1.2):
    """Bildmarke rechts unten. variant: 'rot' (helle Folien) | 'weiss' (rote/dunkle)."""
    from PIL import Image as _Im
    path = LOGOS.get(variant, LOGOS["rot"])
    w_px, h_px = _Im.open(path).size
    w_emu = Inches(size_in); h_emu = int(w_emu * h_px / w_px)
    slide.shapes.add_picture(path, int(GX+GW-w_emu), int(SH-Inches(0.5)-h_emu), w_emu, h_emu)


# ── Kompositions-Komponenten (Leer-Layout) ────────────────────────────────────

def kpi_grid(prs, kpis, kicker="", headline="", folio="", source=""):
    """KPI-Grid mit 3-4 Kennzahlen plakativ.

    Args:
        kpis: list[tuple] mit:
            (zahl, label) — 2-Tupel: Basis-KPI
            (zahl, label, fn) — 3-Tupel: zusätzlich Footnote-Nummer als Superscript
                                          neben der Zahl (für footnote_block-Verweis)

    Beispiel mit Footnotes:
        kpi_grid(prs, [
            ("CHF 2.4 Mio.", "EBIT 2025", 1),   # ¹ wird neben "CHF 2.4 Mio." gezeigt
            ("200",          "Mitarbeitende", 2),
            ("5",            "Verlagsobjekte"),  # ohne Footnote
        ], kicker="Galledia heute", headline="Solide Basis 2025")
        # Danach: footnote_block(s, [(1, "Reporting Q4 2025"), (2, "Stand 31.12.2025")])
    """
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    n=len(kpis); g=Inches(0.5); cw=(GW-(n-1)*g)/n; cy,ch=Inches(3.6),Inches(4.0)
    for i, kpi in enumerate(kpis):
        # Tuple-Form flexibel
        if isinstance(kpi, (tuple, list)) and len(kpi) >= 3:
            zahl, label, fn = kpi[0], kpi[1], kpi[2]
        else:
            zahl, label = kpi[0], kpi[1]; fn = None

        # Footnote-Marker direkt an Label-Ende angehängt (Best Practice statt eigene Shape)
        if fn is not None:
            label = f"{label}{_sup(fn)}"

        x=GX+i*(cw+g); fill=RED if i==0 else GL
        fg=WHITE if i==0 else BLACK; fg2=WHITE if i==0 else G2
        card(s,x,cy,cw,ch,fill)
        tf=tb(s,x+Inches(0.3),cy,cw-Inches(0.6),ch,MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0],zahl, 56,SB,  fg, align=PP_ALIGN.CENTER,space_after=8)
        run(tf.add_paragraph(),label,17,BODY,fg2,align=PP_ALIGN.CENTER,space_after=0)

    _set_footer(s,prs,folio,source); return s

def two_column(prs,head_l,items_l,head_r,items_r,
               col2_red=True,kicker="",headline="",folio="",source=""):
    # ── RIEGEL: Bei zu wenig Inhalt automatisch einspaltig rendern ──────────────
    # Zweispaltig nur bei echtem Vergleich mit >=3 Punkten PRO Spalte.
    # Sonst halbleere Karten → stattdessen dichte einspaltige add_content-Folie.
    if min(len(items_l), len(items_r)) < 3:
        body = head_l + "\n" + "\n".join(f"· {i}" for i in items_l) \
             + "\n\n" + head_r + "\n" + "\n".join(f"· {i}" for i in items_r)
        return add_content(prs, "viel", kicker or "", headline or "", body,
                           folio=folio, source=source)
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    # Kartenhöhe an Inhalt anpassen (Header + Bullets), beide Karten gleich hoch
    rows=max(len(items_l),len(items_r))
    ch=Inches(min(6.8, max(2.4, 1.5 + 0.62*rows)))
    cy=Inches(3.2); cw=(GW-Inches(0.6))/2
    for x,fill,head,fg,items in [
        (GX, GL, head_l, G1, items_l),
        (GX+cw+Inches(0.6), RED if col2_red else GL,
         head_r, WHITE if col2_red else G1, items_r)]:
        card(s,x,cy,cw,ch,fill)
        tf=tb(s,x+Inches(0.6),cy+Inches(0.5),cw-Inches(1.2),ch-Inches(1.0))
        run(tf.paragraphs[0],head,28,SB,  fg,space_after=16)
        for item in items: run(tf.add_paragraph(),f"•  {item}",19,BODY,fg,space_after=10)
    _set_footer(s,prs,folio,source); return s

def flow_pipeline(prs,nodes,cache_text="",kicker="",headline="",folio="",source=""):
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    n=len(nodes); g=Inches(0.4); nw=(GW-(n-1)*g)/n; ny,nh=Inches(3.6),Inches(2.6)
    for i,node in enumerate(nodes):
        titel=node[0] if isinstance(node,(list,tuple)) else node
        sub  =node[1] if isinstance(node,(list,tuple)) and len(node)>1 else ""
        fill=RED if i==n-1 else GL; fg=WHITE if i==n-1 else BLACK; fg2=WHITE if i==n-1 else G2
        x=GX+i*(nw+g); card(s,x,ny,nw,nh,fill)
        tf=tb(s,x+Inches(0.2),ny,nw-Inches(0.4),nh,MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0],titel,20,SB,  fg, align=PP_ALIGN.CENTER,space_after=6)
        if sub: run(tf.add_paragraph(),sub,13,BODY,fg2,align=PP_ALIGN.CENTER,space_after=0)
        if i<n-1:
            at=tb(s,x+nw,ny,g,nh,MSO_ANCHOR.MIDDLE)
            run(at.paragraphs[0],"→",22,SB,G2,align=PP_ALIGN.CENTER,space_after=0)
    if cache_text:
        cy2=ny+nh+Inches(0.55); card(s,GX,cy2,GW,Inches(1.4),TURK)
        tf=tb(s,GX+Inches(0.6),cy2,GW-Inches(1.2),Inches(1.4),MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0],"Semantic Caching",18,SB,  WHITE,space_after=4)
        run(tf.add_paragraph(),cache_text,      15,BODY,WHITE,space_after=0)
    _set_footer(s,prs,folio,source); return s

def timeline(prs,phases,kicker="",headline="",folio="",source=""):
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    ly=Inches(5.4); line=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,GX,ly,GW,Pt(4))
    line.fill.solid(); line.fill.fore_color.rgb=G4; line.line.fill.background(); line.shadow.inherit=False
    n=len(phases); seg=GW/n
    for i,(titel,beschr) in enumerate(phases):
        cx=GX+seg*i+seg/2; col=RED if i==0 else G2
        dot=s.shapes.add_shape(MSO_SHAPE.OVAL,int(cx-Inches(0.22)),int(ly-Inches(0.18)),Inches(0.44),Inches(0.44))
        dot.fill.solid(); dot.fill.fore_color.rgb=col; dot.line.color.rgb=WHITE; dot.line.width=Pt(3); dot.shadow.inherit=False
        tf=tb(s,int(cx-seg/2+Inches(0.3)),ly-Inches(2.6),int(seg-Inches(0.6)),Inches(2.3),MSO_ANCHOR.BOTTOM)
        run(tf.paragraphs[0],titel,22,SB,  BLACK,align=PP_ALIGN.CENTER,space_after=0)
        tf2=tb(s,int(cx-seg/2+Inches(0.3)),ly+Inches(0.5),int(seg-Inches(0.6)),Inches(2.0))
        run(tf2.paragraphs[0],beschr,16,BODY,G2,align=PP_ALIGN.CENTER,space_after=0)
    _set_footer(s,prs,folio,source); return s

def numbered_steps(prs,steps,kicker="",headline="",folio="",source=""):
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    n=len(steps); g=Inches(0.5); cw=(GW-(n-1)*g)/n; cy,ch=Inches(3.6),Inches(5.0)
    for i,(titel,beschr) in enumerate(steps):
        x=GX+i*(cw+g); card(s,x,cy,cw,ch,GL)
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,x+Inches(0.6),cy+Inches(0.55),Inches(1.0),Inches(1.0))
        badge.fill.solid(); badge.fill.fore_color.rgb=RED; badge.line.fill.background(); badge.shadow.inherit=False
        bf=badge.text_frame; bf.word_wrap=False
        run(bf.paragraphs[0],str(i+1),32,SB,WHITE,align=PP_ALIGN.CENTER,space_after=0)
        tf=tb(s,x+Inches(0.6),cy+Inches(1.9),cw-Inches(1.2),ch-Inches(2.3))
        run(tf.paragraphs[0],titel, 24,SB,  BLACK,space_after=10)
        run(tf.add_paragraph(),beschr,17,BODY,G2,  space_after=0)
    _set_footer(s,prs,folio,source); return s

def image_bleed(prs,image_path,kicker="",headline="",body_text="",
                folio="",source="",image_side="right"):
    s=_blank(prs); img_w=Inches(9.2)
    img_x=Inches(20)-img_w if image_side=="right" else 0
    txt_x=GX if image_side=="right" else img_w+Inches(0.5)
    s.shapes.add_picture(image_path,img_x,0,img_w,SH)
    if kicker or headline: _set_header(s,kicker,headline)
    if body_text:
        tf=tb(s,txt_x,Inches(2.4),Inches(9.0),Inches(6.0))
        run(tf.paragraphs[0],body_text,19,BODY,G1,space_after=0)
    _set_footer(s,prs,folio,source); return s


# ═════════════════════════════════════════════════════════════════════════════
# v0.1.0 STUBS — MBB-Level Features (Sprint-Implementierung folgt)
# Siehe ROADMAP.md im Repo-Root für Sprint-Plan und Akzeptanzkriterien.
# ═════════════════════════════════════════════════════════════════════════════

class NotYetImplementedError(NotImplementedError):
    """Funktion ist Teil von v0.1.0-Roadmap, aber noch nicht released."""
    pass

def _stub(name, sprint):
    raise NotYetImplementedError(
        f"{name}() ist in Sprint {sprint} der v0.1.0-Roadmap geplant — noch nicht implementiert. "
        f"Siehe ROADMAP.md."
    )

# ── Sprint 1: Footnotes & Annotations (v0.1.0-alpha.2 — IMPLEMENTED) ─────────

_SUPERSCRIPT = {0:"⁰", 1:"¹", 2:"²", 3:"³", 4:"⁴", 5:"⁵", 6:"⁶", 7:"⁷", 8:"⁸", 9:"⁹"}

def _sup(num):
    """Wandelt Zahl in Unicode-Superscript (1→¹, 12→¹², 123→¹²³)."""
    return "".join(_SUPERSCRIPT.get(int(c), c) for c in str(num))

def add_footnote_marker(slide, target_shape, num=1, position="top-right", size_pt=14, color=None):
    """Setzt ¹/²/³ Superscript-Marker neben einer bestehenden Shape (z.B. KPI-Zahl, Bullet, Headline).

    Args:
        slide: Slide-Objekt
        target_shape: Shape-Objekt mit .left/.top/.width/.height (z.B. add_textbox-Rückgabe)
        num: Footnote-Nummer (1-99)
        position: 'top-right' | 'after' | 'below' | 'top-left'
        size_pt: Schriftgrösse (Default 14pt)
        color: RGBColor (Default BLACK)

    Returns:
        TextFrame des erstellten Markers (für weitere Anpassungen)
    """
    if color is None: color = BLACK
    char = _sup(num)

    # Coordinates aus target_shape extrahieren (kann Shape oder TextFrame sein)
    if hasattr(target_shape, "left"):
        tl, tt, tw, th = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    elif hasattr(target_shape, "_parent"):  # TextFrame → parent Shape
        sh = target_shape._parent
        tl, tt, tw, th = sh.left, sh.top, sh.width, sh.height
    else:
        raise ValueError(f"target_shape muss .left/.top/.width/.height haben oder TextFrame sein, ist: {type(target_shape)}")

    pad = Inches(0.04)
    if position == "top-right":
        x, y = tl + tw + pad, tt
    elif position == "after":
        x, y = tl + tw + pad, tt + th // 3
    elif position == "below":
        x, y = tl + tw - Inches(0.3), tt + th
    elif position == "top-left":
        x, y = tl - Inches(0.3), tt
    else:
        x, y = tl + tw + pad, tt

    tf = tb(slide, int(x), int(y), Inches(0.4), Inches(0.3))
    run(tf.paragraphs[0], char, size_pt, SB, color, space_after=0)
    return tf

def footnote_block(slide, footnotes, y_pos=None, separator="    "):
    """Erzeugt Footnote-Block am Folienfuss (über Standard-Fusszeile).

    Args:
        slide: Slide-Objekt
        footnotes: list[tuple[int, str]] — z.B. [(1, "BDZV 2026, n=94"), (2, "Schätzung Q1")]
        y_pos: y-Koordinate in EMU (Default: Inches(10.30), knapp über Fusszeile)
        separator: String zwischen einzelnen Footnotes (Default 4 Spaces)

    Returns:
        TextFrame des Footnote-Blocks oder None bei leeren footnotes

    Layout: einzeilig, links-bündig, Volte Regular 9pt, Grau G2 (#666666).
    Beispiel-Output: "¹ BDZV 2026, n=94    ² Schätzung Q1    ³ Internes Reporting"
    """
    if not footnotes: return None
    if y_pos is None: y_pos = Inches(10.30)

    parts = [f"{_sup(num)} {text}" for num, text in footnotes]
    full_text = separator.join(parts)

    tf = tb(slide, GX, int(y_pos), GW, Inches(0.40), MSO_ANCHOR.TOP)
    run(tf.paragraphs[0], full_text, 9, BODY, G2, space_after=0)
    return tf

def add_callout(slide, x, y, w, h, text, arrow_to=None, color=None, text_color=None, size_pt=14):
    """Highlight-Box mit optionalem Pfeil zu Ziel-Koordinate — für "hier kippt der Trend"-Hinweise.

    Args:
        slide: Slide-Objekt
        x, y, w, h: Position + Grösse der Box in EMU (z.B. Inches(...))
        text: Box-Text (kurz halten, 1-2 Zeilen)
        arrow_to: tuple (x, y) in EMU — Pfeilspitze zeigt dorthin. None = keine Pfeil.
        color: Background-Farbe (Default RED)
        text_color: Schriftfarbe (Default WHITE auf RED, BLACK sonst)
        size_pt: Schriftgrösse (Default 14)

    Returns:
        Card-Shape (für weitere Anpassungen)
    """
    if color is None: color = RED
    if text_color is None: text_color = WHITE if color == RED else BLACK

    # Card mit abgerundeten Ecken
    sh = card(slide, x, y, w, h, color, corner=0.1)

    # Zentrierter Text
    tf = tb(slide, x + Inches(0.15), y, w - Inches(0.3), h, MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], text, size_pt, SB, text_color, align=PP_ALIGN.CENTER, space_after=0)

    # Optional: Pfeil zur Ziel-Koordinate
    if arrow_to:
        tx, ty = arrow_to
        cx_center = x + w // 2
        cy_center = y + h // 2

        # Pfeil-Start an näherer Box-Kante bestimmen
        if tx < x:
            start_x, start_y = x, cy_center
        elif tx > x + w:
            start_x, start_y = x + w, cy_center
        elif ty < y:
            start_x, start_y = cx_center, y
        else:
            start_x, start_y = cx_center, y + h

        from pptx.enum.shapes import MSO_CONNECTOR
        from lxml import etree
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, int(start_x), int(start_y), int(tx), int(ty))
        conn.line.color.rgb = color
        conn.line.width = Pt(2)
        # Pfeilspitze am Ende anhängen
        ln = conn.line._get_or_add_ln()
        tailEnd = ln.find(qn('a:tailEnd'))
        if tailEnd is None:
            tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', 'med')
        tailEnd.set('len', 'med')

    return sh

def add_highlight_line(slide, x1, y1, x2, y2, label="", color=None, dashed=False, label_pos="end", label_color=None):
    """Akzentlinie für Charts/Folien — horizontal, vertikal oder diagonal.

    Args:
        slide: Slide-Objekt
        x1, y1, x2, y2: Linien-Koordinaten in EMU
        label: optionaler Beschriftungstext
        color: Linienfarbe (Default RED)
        dashed: True für gestrichelte Linie
        label_pos: 'end' | 'start' | 'mid' — wo Label platzieren
        label_color: Schriftfarbe des Labels (Default = color)

    Returns:
        Connector-Shape

    Beispiele:
        # Horizontale Trend-Linie über einem Chart bei y=Inches(4.0)
        add_highlight_line(s, Inches(1), Inches(4), Inches(19), Inches(4), "Zielwert 850k", dashed=True)

        # Vertikale Markierung "heute" auf Timeline
        add_highlight_line(s, Inches(10), Inches(3.5), Inches(10), Inches(6.5), "Heute", color=RED)
    """
    if color is None: color = RED
    if label_color is None: label_color = color

    from pptx.enum.shapes import MSO_CONNECTOR
    from lxml import etree
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(2)

    if dashed:
        ln = conn.line._get_or_add_ln()
        prstDash = ln.find(qn('a:prstDash'))
        if prstDash is None:
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')

    if label:
        if label_pos == "end":
            lx, ly = x2 + Inches(0.12), y2 - Inches(0.18)
        elif label_pos == "start":
            lx, ly = x1 - Inches(1.5), y1 - Inches(0.18)
        else:  # mid
            lx = (x1 + x2) // 2 - Inches(0.75)
            ly = (y1 + y2) // 2 - Inches(0.25)
        tf = tb(slide, int(lx), int(ly), Inches(2.5), Inches(0.35))
        run(tf.paragraphs[0], label, 11, SB, label_color, space_after=0)

    return conn

def add_value_label(slide, anchor_xy, value, fmt="+%.1f%%", size_pt=14, color=None, anchor=MSO_ANCHOR.TOP):
    """Formatierte Wert-Annotation an Chart-Punkt — für "+23%"-Marker an Datenpunkten.

    Args:
        slide: Slide-Objekt
        anchor_xy: tuple (x, y) in EMU — Position des Labels
        value: numerischer Wert (int oder float)
        fmt: Python format-string (Default '+%.1f%%' für '+23.4%')
             Beispiele: '%.0f' für '23', 'CHF %.1fM' für 'CHF 2.4M', '%+d' für '+23'
        size_pt: Schriftgrösse (Default 14)
        color: Schriftfarbe (Default RED)
        anchor: vertical anchor (MSO_ANCHOR.TOP/MIDDLE/BOTTOM)

    Returns:
        TextFrame des Labels
    """
    if color is None: color = RED

    try:
        text = fmt % value
    except (TypeError, ValueError):
        text = str(value)

    x, y = anchor_xy
    tf = tb(slide, int(x), int(y), Inches(1.5), Inches(0.4), anchor)
    run(tf.paragraphs[0], text, size_pt, SB, color, space_after=0)
    return tf

# ── Sprint 2: Shape-based Data Visualization (v0.1.0-alpha.3 — IMPLEMENTED) ──

# Chart-Bereich für Sprint-2/3/4/5/6-Layouts (gemeinsam genutzt)
# v0.1.3: exakte Inhaltszone laut User-Vorgabe (Höhe 15.9 cm, Breite 45.86 cm,
#         beginnend bei Horizontal 2.49 cm und Vertikal 7.49 cm).
#         → top=2.95", height=6.26", endet bei y=9.21"
_CHART_TOP    = Inches(2.95)   # = 7.49 cm
_CHART_HEIGHT = Inches(6.26)   # = 15.9 cm — endet bei y=9.21"
_CHART_LEFT   = GX             # 0.979" = 2.49 cm
_CHART_RIGHT  = GX + GW        # ~19.005"
_CHART_WIDTH  = GW             # 18.026" = 45.79 cm

def _fmt_num(val, with_sign=False):
    """Formatiert Zahl mit Schweizer Tausendertrenner ('). +/- Zeichen optional."""
    if with_sign:
        return f"{val:+,}".replace(",", "'")
    return f"{val:,}".replace(",", "'")

def waterfall(prs, segments, kicker="", headline="", folio="", source=""):
    """Brücken-Analyse (z.B. EBIT Vorjahr → Treiber → EBIT aktuell).

    Args:
        prs: Presentation-Objekt
        segments: list[tuple[label, value, kind]]
            kind in {"absolute", "delta_pos", "delta_neg", "total"}
            - "absolute": Start- oder Pivot-Wert von 0 (grau G2)
            - "delta_pos": Positive Veränderung (türkis TURK)
            - "delta_neg": Negative Veränderung (rot RED)
            - "total":    Zwischen-/Endsumme von 0 (dunkelgrau G1)
        kicker, headline, folio, source: wie sonst

    Beispiel:
        waterfall(prs, [
            ("EBIT 2025",      2400, "absolute"),
            ("Umsatz +",        380, "delta_pos"),
            ("Kosten +",       -180, "delta_neg"),
            ("Synergien +",     120, "delta_pos"),
            ("EBIT 2026",      2720, "total"),
        ], kicker="Ergebnis-Bridge", headline="EBIT-Treiber 2025–2026", folio="3")
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    n = len(segments)
    gap = Inches(0.20)
    bar_w = (_CHART_WIDTH - (n-1)*gap) / n

    # 1. Berechne from_val/to_val pro Bar im DATA-Space
    running = 0
    bars = []  # (from_v, to_v, kind, label, val)
    for label, val, kind in segments:
        if kind == "absolute":
            from_v, to_v = 0, val; running = val
        elif kind == "total":
            from_v, to_v = 0, val; running = val
        elif kind == "delta_pos":
            from_v, to_v = running, running + val; running += val
        elif kind == "delta_neg":
            from_v, to_v = running, running + val; running += val
        else:
            raise ValueError(f"Unbekannter Segment-Typ: {kind}")
        bars.append((from_v, to_v, kind, label, val))

    # 2. Y-Achsen-Skala bestimmen (alle Werte berücksichtigen, von 0 ab)
    max_v = max(max(abs(b[0]), abs(b[1])) for b in bars)
    if max_v == 0: max_v = 1
    scale = float(_CHART_HEIGHT) / max_v
    chart_baseline = int(_CHART_TOP) + int(_CHART_HEIGHT)  # y bei value=0

    # 3. Bars + Labels + Connectors zeichnen
    from pptx.enum.shapes import MSO_CONNECTOR
    from lxml import etree
    bar_tops = []  # y-coord des top-edges (für connector zur nächsten bar)
    for i, (fv, tv, kind, label, val) in enumerate(bars):
        x = _CHART_LEFT + i * (bar_w + gap)
        y_top = chart_baseline - int(max(fv, tv) * scale)
        y_bot = chart_baseline - int(min(fv, tv) * scale)
        bar_h = max(y_bot - y_top, Inches(0.06))  # min. 5px

        if kind == "delta_pos":   color = TURK
        elif kind == "delta_neg": color = RED
        elif kind == "total":     color = G1
        else:                     color = G2  # absolute

        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y_top), int(bar_w), int(bar_h))
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        sh.line.fill.background(); sh.shadow.inherit = False

        # Wert-Label oberhalb der Bar
        val_text = _fmt_num(val, with_sign=(kind in ("delta_pos", "delta_neg")))
        vtf = tb(s, int(x), int(y_top - Inches(0.42)), int(bar_w), Inches(0.36))
        run(vtf.paragraphs[0], val_text, 14, SB, color, align=PP_ALIGN.CENTER, space_after=0)

        # Kategorie-Label unter Baseline
        ltf = tb(s, int(x), chart_baseline + int(Inches(0.10)), int(bar_w), Inches(0.55))
        run(ltf.paragraphs[0], label, 12, BODY, BLACK, align=PP_ALIGN.CENTER, space_after=0)

        bar_tops.append(y_top)

    # Connector-Linien zwischen aufeinanderfolgenden Bars (top-edge to top-edge)
    for i in range(n - 1):
        # Connector zeichnen wir vom oberen Ende der current bar zur next bar
        # nur sinnvoll, wenn nächste Bar ein delta ist (sitzt direkt am running-total)
        if bars[i+1][2] in ("delta_pos", "delta_neg"):
            conn_y = chart_baseline - int(bars[i+1][0] * scale)  # next bar's from_val
            x1 = _CHART_LEFT + i * (bar_w + gap) + bar_w
            x2 = _CHART_LEFT + (i+1) * (bar_w + gap)
            line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(conn_y), int(x2), int(conn_y))
            line.line.color.rgb = G3; line.line.width = Pt(1)
            ln = line.line._get_or_add_ln()
            prstDash = ln.find(qn('a:prstDash'))
            if prstDash is None:
                prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'dash')

    _set_footer(s, prs, folio, source); return s

def matrix_2x2(prs, x_axis, y_axis, items, kicker="", headline="", folio="", source=""):
    """2×2-Matrix (Impact/Effort, Risiko/Chance, BCG-light, Priorisierung).

    Args:
        x_axis: tuple (label_low, label_high) — z.B. ("Aufwand niedrig", "Aufwand hoch")
        y_axis: tuple (label_low, label_high) — z.B. ("Impact niedrig", "Impact hoch")
        items: list[tuple[x, y, label]] — x,y in [0,1] (0=low, 1=high), label = Item-Name
            Optional: tuple[x, y, label, size] mit size in [0.5, 2.0] für Bubble-Grösse

    Quadranten:
        Top-Right (x>0.5, y>0.5):    RED (Quick Win / Star)
        Top-Left  (x<0.5, y>0.5):    TURK (Strategisch wertvoll)
        Bottom-Right (x>0.5, y<0.5): G3 (Vermeiden)
        Bottom-Left (x<0.5, y<0.5):  G4 (Aufschieben)
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    # Quadranten-Background (hellere Töne, damit Items lesbar)
    qx, qy = _CHART_LEFT + Inches(1.5), _CHART_TOP
    qw = _CHART_WIDTH - Inches(2.5)  # Platz für Y-Label links
    qh = _CHART_HEIGHT - Inches(0.5)  # Platz für X-Label unten
    half_w, half_h = qw // 2, qh // 2

    # Quadranten als 4 farbige Rechtecke (helle Töne via Transparenz-Simulation)
    quads = [
        (qx,           qy,           half_w, half_h, GL),    # top-left (light)
        (qx + half_w,  qy,           half_w, half_h, RGBColor(0xFD, 0xE7, 0xED)),  # top-right (light red)
        (qx,           qy + half_h,  half_w, half_h, GL),    # bottom-left
        (qx + half_w,  qy + half_h,  half_w, half_h, GL),    # bottom-right
    ]
    for x_, y_, w_, h_, fill in quads:
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x_), int(y_), int(w_), int(h_))
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = G4; rect.line.width = Pt(0.75)
        rect.shadow.inherit = False

    # Achsen-Labels
    # X-Achse: links/rechts unter der Box
    ax_l = tb(s, int(qx), int(qy + qh + Inches(0.12)), int(half_w), Inches(0.35))
    run(ax_l.paragraphs[0], x_axis[0], 12, BODY, G2, align=PP_ALIGN.LEFT, space_after=0)
    ax_r = tb(s, int(qx + half_w), int(qy + qh + Inches(0.12)), int(half_w), Inches(0.35))
    run(ax_r.paragraphs[0], x_axis[1], 12, SB, BLACK, align=PP_ALIGN.RIGHT, space_after=0)
    # Y-Achse: oberhalb/unterhalb links der Box (vertical text wäre besser, hier horizontal vor Quadrant)
    ay_h = tb(s, int(qx - Inches(1.4)), int(qy + Inches(0.05)), Inches(1.3), Inches(0.35))
    run(ay_h.paragraphs[0], y_axis[1], 12, SB, BLACK, align=PP_ALIGN.RIGHT, space_after=0)
    ay_l = tb(s, int(qx - Inches(1.4)), int(qy + qh - Inches(0.4)), Inches(1.3), Inches(0.35))
    run(ay_l.paragraphs[0], y_axis[0], 12, BODY, G2, align=PP_ALIGN.RIGHT, space_after=0)

    # Items als Bubbles
    for item in items:
        if len(item) == 4:
            ix, iy, label, size_mult = item
        else:
            ix, iy, label = item; size_mult = 1.0
        bubble_d = Inches(0.50 * size_mult)
        # Pixel-Koordinaten: ix=0 → left edge of quadrant box; ix=1 → right edge
        px = qx + int(ix * qw) - bubble_d // 2
        # y in pptx grows down; iy=0 = bottom, iy=1 = top → flip
        py = qy + int((1 - iy) * qh) - bubble_d // 2

        # Bubble (filled circle)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, int(px), int(py), int(bubble_d), int(bubble_d))
        circle.fill.solid(); circle.fill.fore_color.rgb = RED
        circle.line.color.rgb = WHITE; circle.line.width = Pt(2)
        circle.shadow.inherit = False

        # Label rechts neben Bubble
        ltf = tb(s, int(px + bubble_d + Inches(0.05)), int(py - Inches(0.05)), Inches(3.5), Inches(0.5))
        run(ltf.paragraphs[0], label, 12, SB, BLACK, align=PP_ALIGN.LEFT, space_after=0)

    _set_footer(s, prs, folio, source); return s

def mekko(prs, segments, kicker="", headline="", folio="", source=""):
    """Marimekko-Chart: Breite = Marktgrösse, Höhe = Anteile pro Markt.

    Args:
        segments: list[dict] mit:
            "name":   str (Markt-Bezeichnung)
            "size":   float in [0,1] (relativer Marktanteil — Summe sollte ≈1 sein)
            "shares": list[tuple[player_name, share]] mit share in [0,1]
                      (erster Eintrag = Galledia = RED, Rest in Grautönen)

    Beispiel:
        mekko(prs, segments=[
            {"name": "Print", "size": 0.45, "shares": [("Galledia", 0.35), ("Wettbewerb A", 0.25), ("Rest", 0.40)]},
            {"name": "Digital", "size": 0.30, "shares": [("Galledia", 0.20), ("Wettbewerb A", 0.45), ("Rest", 0.35)]},
            {"name": "Event", "size": 0.25, "shares": [("Galledia", 0.40), ("Wettbewerb A", 0.30), ("Rest", 0.30)]},
        ], kicker="Marktanteile", headline="Galledia stärkster in Event")
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    total_size = sum(seg["size"] for seg in segments)
    if total_size == 0: total_size = 1

    # Player-Farben (erste = RED für Galledia, dann Grautöne abnehmend dunkel)
    PLAYER_COLORS = [RED, G1, G2, G3, G4, TURK]

    # Header-Block IN der Inhaltszone (oberhalb der Flächen), nicht in der Schutzzone
    HEADER_H = Inches(0.65)              # Höhe Title + %
    flaechen_top = _CHART_TOP + HEADER_H  # Flächen starten nach Header
    flaechen_h = _CHART_HEIGHT - HEADER_H  # = 6.26 - 0.65 = 5.61"

    x_cursor = _CHART_LEFT
    for seg in segments:
        col_w = int(_CHART_WIDTH * (seg["size"] / total_size))

        # Title direkt am oberen Rand der Inhaltszone
        head_tf = tb(s, int(x_cursor), int(_CHART_TOP), col_w, Inches(0.32))
        run(head_tf.paragraphs[0], seg["name"], 14, SB, BLACK, align=PP_ALIGN.CENTER, space_after=0)
        # % direkt darunter
        size_tf = tb(s, int(x_cursor), int(_CHART_TOP + Inches(0.34)), col_w, Inches(0.24))
        run(size_tf.paragraphs[0], f"{seg['size']*100:.0f}%", 11, BODY, G2, align=PP_ALIGN.CENTER, space_after=0)

        # Shares stacken, beginnen unter dem Header
        y_cursor = int(flaechen_top)
        shares_norm_total = sum(sh[1] for sh in seg["shares"]) or 1
        for j, (player, share) in enumerate(seg["shares"]):
            seg_h = int(flaechen_h * (share / shares_norm_total))
            color = PLAYER_COLORS[j % len(PLAYER_COLORS)]
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x_cursor), y_cursor, col_w, seg_h)
            rect.fill.solid(); rect.fill.fore_color.rgb = color
            rect.line.color.rgb = WHITE; rect.line.width = Pt(1.5)
            rect.shadow.inherit = False

            # Player-Label + Anteil im Rechteck (nur wenn gross genug)
            if seg_h > Inches(0.45):
                text_color = WHITE if color in (RED, G1, G2) else BLACK
                lbl = f"{player}\n{share*100:.0f}%"
                ltf = tb(s, int(x_cursor + Inches(0.08)), y_cursor, col_w - int(Inches(0.16)), seg_h, MSO_ANCHOR.MIDDLE)
                run(ltf.paragraphs[0], player, 11, SB, text_color, space_after=2)
                run(ltf.add_paragraph(), f"{share*100:.0f}%", 14, SB, text_color, space_after=0)
            y_cursor += seg_h

        x_cursor += col_w + Inches(0.05)  # kleine Spalten-Trennung

    _set_footer(s, prs, folio, source); return s

def stacked_bar_pct(prs, categories, series, kicker="", headline="", folio="", source=""):
    """100% Stacked Bar — Vergleich Anteils-Entwicklung über Kategorien.

    Args:
        categories: list[str] — z.B. ["2024", "2025", "2026"]
        series: list[tuple[name, values]] — values pro Kategorie, in der Reihenfolge der categories
                Beispiel: [("Print", [45, 38, 30]), ("Digital", [35, 42, 48]), ("Event", [20, 20, 22])]
                Werte können absolute Zahlen sein — werden zu 100% normalisiert.

    Layout:
        - Vertical Bars (eine pro Kategorie), nebeneinander
        - Bars zu 100% gestackt
        - Legend rechts neben Chart
        - Galledia-Rot für erste Serie, dann Grautöne
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    n_cat = len(categories)
    n_ser = len(series)

    legend_w = Inches(3.0)
    chart_w = _CHART_WIDTH - legend_w
    gap = Inches(0.30)
    bar_w = (chart_w - (n_cat - 1) * gap) / n_cat

    SERIES_COLORS = [RED, G1, G2, G3, TURK, G4]

    # Pro Kategorie: Normalisierung und Stacking
    for ci, cat in enumerate(categories):
        col_total = sum(ser[1][ci] for ser in series) or 1
        x = _CHART_LEFT + ci * (bar_w + gap)
        y_cursor = int(_CHART_TOP)
        for si, (name, values) in enumerate(series):
            val = values[ci]
            pct = val / col_total
            seg_h = int(_CHART_HEIGHT * pct)
            color = SERIES_COLORS[si % len(SERIES_COLORS)]

            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), y_cursor, int(bar_w), seg_h)
            rect.fill.solid(); rect.fill.fore_color.rgb = color
            rect.line.color.rgb = WHITE; rect.line.width = Pt(1.5)
            rect.shadow.inherit = False

            # Prozent-Label im Segment (nur wenn gross genug)
            if seg_h > Inches(0.35):
                text_color = WHITE if color in (RED, G1, G2) else BLACK
                ltf = tb(s, int(x), y_cursor, int(bar_w), seg_h, MSO_ANCHOR.MIDDLE)
                run(ltf.paragraphs[0], f"{pct*100:.0f}%", 14, SB, text_color, align=PP_ALIGN.CENTER, space_after=0)
            y_cursor += seg_h

        # Kategorie-Label unter Bar
        cat_tf = tb(s, int(x), int(_CHART_TOP + _CHART_HEIGHT + Inches(0.10)), int(bar_w), Inches(0.40))
        run(cat_tf.paragraphs[0], cat, 13, SB, BLACK, align=PP_ALIGN.CENTER, space_after=0)

    # Legende rechts neben Chart
    legend_x = _CHART_LEFT + chart_w + Inches(0.40)
    legend_y = _CHART_TOP + Inches(0.50)
    for si, (name, _) in enumerate(series):
        color = SERIES_COLORS[si % len(SERIES_COLORS)]
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(legend_x), int(legend_y + si * Inches(0.45)),
                                Inches(0.30), Inches(0.30))
        sq.fill.solid(); sq.fill.fore_color.rgb = color
        sq.line.fill.background(); sq.shadow.inherit = False
        ltf = tb(s, int(legend_x + Inches(0.40)), int(legend_y + si * Inches(0.45) - Inches(0.04)),
                 Inches(2.5), Inches(0.40))
        run(ltf.paragraphs[0], name, 13, BODY, BLACK, space_after=0)

    _set_footer(s, prs, folio, source); return s

def quadrant_map(prs, items_xy, axes_labels, kicker="", headline="", folio="", source=""):
    """Strategische Positionierung — Items als Scatter mit beschrifteten Achsen.

    Unterscheidet sich von matrix_2x2 durch:
    - Keine harten Quadranten-Farb-Boxen, stattdessen Achsenkreuz
    - Mehr Platz für lange Labels
    - Optional: hervorgehobener "Sweet Spot"-Bereich

    Args:
        items_xy: list[tuple[x, y, label]] mit x,y in [0,1]
                  Optional: tuple[x, y, label, is_us] mit is_us=True für eigene Positionierung (Galledia)
        axes_labels: tuple (x_label, y_label) — z.B. ("Reichweite", "Tiefe der Beziehung")

    Beispiel:
        quadrant_map(prs,
            items_xy=[
                (0.7, 0.8, "Galledia Fachmedien", True),
                (0.3, 0.6, "Wettbewerb A"),
                (0.5, 0.4, "Wettbewerb B"),
                (0.2, 0.3, "Wettbewerb C"),
            ],
            axes_labels=("Reichweite", "Tiefe der Beziehung"),
            kicker="Wettbewerbs-Landschaft",
            headline="Galledia stark positioniert")
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    qx, qy = _CHART_LEFT + Inches(1.2), _CHART_TOP
    qw = _CHART_WIDTH - Inches(2.0)
    qh = _CHART_HEIGHT - Inches(0.5)

    # Hintergrund (leicht graustichig)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(qx), int(qy), int(qw), int(qh))
    bg.fill.solid(); bg.fill.fore_color.rgb = GL
    bg.line.color.rgb = G4; bg.line.width = Pt(0.75)
    bg.shadow.inherit = False

    # Achsen-Kreuz (Mittel-Linien)
    from pptx.enum.shapes import MSO_CONNECTOR
    mid_x, mid_y = qx + qw // 2, qy + qh // 2
    # Horizontal mid-line
    line_h = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(qx), int(mid_y), int(qx + qw), int(mid_y))
    line_h.line.color.rgb = G3; line_h.line.width = Pt(0.75)
    # Vertical mid-line
    line_v = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(mid_x), int(qy), int(mid_x), int(qy + qh))
    line_v.line.color.rgb = G3; line_v.line.width = Pt(0.75)

    # Achsen-Labels
    x_lbl = tb(s, int(qx), int(qy + qh + Inches(0.10)), int(qw), Inches(0.30))
    run(x_lbl.paragraphs[0], f"→ {axes_labels[0]}", 12, SB, BLACK, align=PP_ALIGN.CENTER, space_after=0)
    y_lbl = tb(s, int(qx - Inches(1.15)), int(qy + qh // 2 - Inches(0.18)), Inches(1.05), Inches(0.35))
    run(y_lbl.paragraphs[0], f"↑ {axes_labels[1]}", 12, SB, BLACK, align=PP_ALIGN.RIGHT, space_after=0)

    # Items als Bubbles — eigene Position (is_us=True) in RED, Wettbewerb in G2
    for item in items_xy:
        if len(item) == 4:
            ix, iy, label, is_us = item
        else:
            ix, iy, label = item; is_us = False
        bubble_d = Inches(0.60 if is_us else 0.45)
        px = qx + int(ix * qw) - bubble_d // 2
        py = qy + int((1 - iy) * qh) - bubble_d // 2

        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, int(px), int(py), int(bubble_d), int(bubble_d))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RED if is_us else G2
        circle.line.color.rgb = WHITE; circle.line.width = Pt(2)
        circle.shadow.inherit = False

        # Label — Position abhängig von Lage in Box (vermeide Rand-Überlauf)
        label_x = px + bubble_d + Inches(0.06)
        label_w = Inches(3.0)
        # Wenn Bubble nah am rechten Rand: Label links
        if ix > 0.75:
            label_x = px - label_w - Inches(0.06)
            align = PP_ALIGN.RIGHT
        else:
            align = PP_ALIGN.LEFT
        ltf = tb(s, int(label_x), int(py + Inches(0.04)), int(label_w), Inches(0.4))
        run(ltf.paragraphs[0], label, 12,
            SB if is_us else BODY,
            BLACK,
            align=align, space_after=0)

    _set_footer(s, prs, folio, source); return s

# ── Sprint 3: matplotlib-Charts (v0.1.0-alpha.4 — IMPLEMENTED) ───────────────

# Galledia-CI Palette für matplotlib (in Hex statt RGBColor)
_GALLEDIA_HEX = {
    "RED":   "#E61C52",
    "G1":    "#404040",
    "G2":    "#666666",
    "G3":    "#A6A6A6",
    "G4":    "#D9D9D9",
    "GL":    "#F2F2F5",
    "TURK":  "#22AA9F",
    "BLACK": "#000000",
    "WHITE": "#FFFFFF",
}
_GALLEDIA_PALETTE = ["#E61C52", "#22AA9F", "#404040", "#666666", "#A6A6A6", "#D9D9D9"]
_MATPLOTLIB_INITIALIZED = False

def _init_matplotlib():
    """Lazy-load matplotlib, registriert Volte-Fonts + setzt Galledia-Style.
    Wird beim ersten Chart-Aufruf einmal ausgeführt."""
    global _MATPLOTLIB_INITIALIZED
    try:
        import matplotlib
        matplotlib.use('Agg')  # headless, kein Display nötig
        import matplotlib.pyplot as plt
        from matplotlib import font_manager
    except ImportError:
        raise RuntimeError(
            "matplotlib nicht installiert. Sprint-3-Charts brauchen: "
            "pip install matplotlib --break-system-packages"
        )

    if _MATPLOTLIB_INITIALIZED:
        return plt

    # Volte-Fonts aus assets/fonts/ registrieren
    font_dir = os.path.join(os.path.dirname(__file__), "assets", "fonts")
    volte_found = False
    if os.path.isdir(font_dir):
        for f in os.listdir(font_dir):
            if f.endswith(".otf"):
                try:
                    font_manager.fontManager.addfont(os.path.join(font_dir, f))
                    if "Volte" in f:
                        volte_found = True
                except Exception:
                    pass

    # Galledia-CI Defaults
    primary_font = "Volte" if volte_found else "DejaVu Sans"
    plt.rcParams.update({
        "font.family": primary_font,
        "font.size": 11,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "axes.grid": True,
        "grid.alpha": 0.25,
        "grid.color": "#A6A6A6",
        "grid.linestyle": "-",
        "grid.linewidth": 0.6,
        "axes.edgecolor": "#666666",
        "axes.labelcolor": "#000000",
        "axes.titlecolor": "#000000",
        "xtick.color": "#666666",
        "ytick.color": "#666666",
        "figure.facecolor": "white",
        "axes.facecolor": "white",
        "savefig.dpi": 150,
        "savefig.facecolor": "white",
        "axes.prop_cycle": plt.cycler("color", _GALLEDIA_PALETTE),
    })
    _MATPLOTLIB_INITIALIZED = True
    return plt

def _render_chart_png(fig, prefix="galledia_chart_"):
    """Speichert matplotlib-Figure als PNG (tempfile) und gibt Pfad zurück."""
    import tempfile
    f = tempfile.NamedTemporaryFile(suffix=".png", delete=False, prefix=prefix)
    f.close()
    fig.savefig(f.name, dpi=150, bbox_inches="tight", pad_inches=0.15, facecolor="white")
    return f.name

def _embed_chart(prs, png_path, kicker, headline, folio, source):
    """Erstellt Slide, setzt Header/Footer, embedded das Chart-PNG in den Chart-Bereich."""
    s = _blank(prs)
    if kicker or headline:
        _set_header(s, kicker, headline)
    s.shapes.add_picture(png_path, int(_CHART_LEFT), int(_CHART_TOP),
                          width=int(_CHART_WIDTH), height=int(_CHART_HEIGHT))  # Aspect 4.006 — matched zu figsize
    _set_footer(s, prs, folio, source)
    return s

def chart_bar(prs, data, kicker="", headline="", folio="", source="",
              y_label="", value_fmt="{:,.0f}", horizontal=False, color=None):
    """Bar-Chart via matplotlib in Galledia-CI-Farben.

    Args:
        prs: Presentation-Objekt
        data: list[tuple[label, value]] ODER dict[label, value]
        y_label: optional y-Achsen-Beschriftung (bei horizontal → x-Beschriftung)
        value_fmt: Format-String für Wert-Labels (Default '{:,.0f}' → 2,400)
        horizontal: True für horizontale Bars
        color: einzelne Bar-Farbe (Default Galledia-Rot)

    Beispiel:
        chart_bar(prs, [
            ("Print", 1620), ("Digital", 980), ("Event", 450), ("Sonstige", 220),
        ], kicker="Umsatz 2026", headline="Sparten-Übersicht", y_label="CHF k")
    """
    plt = _init_matplotlib()
    if isinstance(data, dict):
        data = list(data.items())
    labels = [str(d[0]) for d in data]
    values = [d[1] for d in data]
    color = color or _GALLEDIA_HEX["RED"]

    fig, ax = plt.subplots(figsize=(18, 6.26))
    if horizontal:
        bars = ax.barh(labels, values, color=color, edgecolor="white", linewidth=1)
        for bar, v in zip(bars, values):
            ax.text(v, bar.get_y() + bar.get_height()/2, " " + value_fmt.format(v).replace(",", "'"),
                    va="center", fontsize=11, color="#000000")
        if y_label: ax.set_xlabel(y_label)
        ax.invert_yaxis()  # erste Bar oben
    else:
        bars = ax.bar(labels, values, color=color, edgecolor="white", linewidth=1)
        for bar, v in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width()/2, v, value_fmt.format(v).replace(",", "'"),
                    ha="center", va="bottom", fontsize=11, color="#000000")
        if y_label: ax.set_ylabel(y_label)

    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_chart_bar_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

def chart_line(prs, x, series, kicker="", headline="", folio="", source="",
               annotations=None, y_label="", x_label="", show_markers=True):
    """Line-Chart mit mehreren Serien und optionalen Annotations.

    Args:
        x: list — x-Werte (numerisch oder strings)
        series: dict[name, list[values]] — z.B. {"Galledia": [100, 108, 124], "Markt": [100, 105, 110]}
        annotations: list[dict] mit "x", "y", "text" — Highlight-Punkte mit Beschriftung
        show_markers: True zeigt Datenpunkte als Kreise

    Beispiel:
        chart_line(prs,
            x=["2024", "2025", "2026"],
            series={"Galledia": [100, 108, 124], "Markt": [100, 105, 110]},
            kicker="Indexierte Entwicklung", headline="Galledia wächst überproportional",
            annotations=[{"x": "2026", "y": 124, "text": "+24% vs. Markt"}])
    """
    plt = _init_matplotlib()
    fig, ax = plt.subplots(figsize=(18, 6.26))
    for i, (name, values) in enumerate(series.items()):
        color = _GALLEDIA_PALETTE[i % len(_GALLEDIA_PALETTE)]
        ax.plot(x, values, label=name, color=color, linewidth=2.5,
                marker="o" if show_markers else None, markersize=7,
                markeredgecolor="white", markeredgewidth=1.5)
    ax.legend(loc="best", frameon=False, fontsize=11)
    if y_label: ax.set_ylabel(y_label)
    if x_label: ax.set_xlabel(x_label)

    # Annotations mit Pfeilen
    if annotations:
        for ann in annotations:
            ax.annotate(ann["text"], xy=(ann["x"], ann["y"]),
                        xytext=(15, 15), textcoords="offset points",
                        fontsize=11, color=_GALLEDIA_HEX["RED"], fontweight="bold",
                        arrowprops=dict(arrowstyle="->", color=_GALLEDIA_HEX["RED"], lw=1.5))

    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_chart_line_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

def chart_stacked_bar(prs, categories, series, kicker="", headline="", folio="", source="",
                       y_label="", value_fmt="{:,.0f}"):
    """Stacked Bar (ABSOLUTE Werte) — für `100% Stacking` siehe stacked_bar_pct.

    Args:
        categories: list[str] — z.B. ["2024", "2025", "2026"]
        series: list[tuple[name, values]] — Werte pro Kategorie
            Beispiel: [("Print", [1620, 1450, 1280]), ("Digital", [620, 780, 980]), ...]

    Beispiel:
        chart_stacked_bar(prs,
            categories=["2024", "2025", "2026"],
            series=[("Print", [1620, 1450, 1280]),
                    ("Digital", [620, 780, 980]),
                    ("Event", [380, 410, 450])],
            kicker="Umsatz nach Sparte", headline="Digital kompensiert Print", y_label="CHF k")
    """
    plt = _init_matplotlib()
    fig, ax = plt.subplots(figsize=(18, 6.26))

    n_cat = len(categories)
    bottom = [0] * n_cat
    for i, (name, values) in enumerate(series):
        color = _GALLEDIA_PALETTE[i % len(_GALLEDIA_PALETTE)]
        bars = ax.bar(categories, values, label=name, bottom=bottom,
                      color=color, edgecolor="white", linewidth=1)
        # Wert-Labels in Segmenten (wenn gross genug)
        for j, bar in enumerate(bars):
            h = bar.get_height()
            if h > max([sum(s[1]) for s in series]) * 0.05:  # nur grosse Segmente labeln
                ax.text(bar.get_x() + bar.get_width()/2,
                        bar.get_y() + h/2,
                        value_fmt.format(values[j]).replace(",", "'"),
                        ha="center", va="center", fontsize=10,
                        color="white" if color in (_GALLEDIA_HEX["RED"], _GALLEDIA_HEX["G1"], _GALLEDIA_HEX["G2"]) else "#000")
        bottom = [b + v for b, v in zip(bottom, values)]

    ax.legend(loc="upper right", frameon=False, fontsize=11)
    if y_label: ax.set_ylabel(y_label)
    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_chart_stacked_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

def chart_histogram(prs, values, bins=10, kicker="", headline="", folio="", source="",
                     y_label="Häufigkeit", x_label="", color=None):
    """Verteilungs-Histogramm (z.B. Antwortzeit-Verteilung, Lead-Score-Verteilung).

    Args:
        values: list[float] — die zu binnenden Werte
        bins: Anzahl Bins (int) ODER explizite Bin-Kanten (list)

    Beispiel:
        chart_histogram(prs, lead_scores, bins=20,
            kicker="Lead-Qualität", headline="Score-Verteilung Mai 2026",
            x_label="Score")
    """
    plt = _init_matplotlib()
    color = color or _GALLEDIA_HEX["RED"]
    fig, ax = plt.subplots(figsize=(18, 6.26))
    ax.hist(values, bins=bins, color=color, edgecolor="white", linewidth=1)
    if y_label: ax.set_ylabel(y_label)
    if x_label: ax.set_xlabel(x_label)
    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_chart_hist_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

def chart_scatter(prs, points, kicker="", headline="", folio="", source="",
                   x_label="", y_label="", labels=True, point_size=80):
    """Scatter-Plot mit Label-Annotations.

    Args:
        points: list[tuple[x, y]] ODER list[tuple[x, y, label]]
            Optional 4. Element: tuple[x, y, label, color_hex]
        labels: True zeigt Label neben jedem Punkt
        point_size: matplotlib s-Parameter für Punkt-Grösse

    Beispiel:
        chart_scatter(prs,
            points=[(0.7, 0.8, "Galledia"), (0.3, 0.6, "Wettbewerb A"), (0.5, 0.4, "Wettbewerb B")],
            x_label="Reichweite", y_label="Tiefe",
            kicker="Positionierung", headline="Galledia ganz oben rechts")
    """
    plt = _init_matplotlib()
    fig, ax = plt.subplots(figsize=(18, 6.26))

    for i, pt in enumerate(points):
        if len(pt) >= 4:
            x_, y_, lbl, col = pt[0], pt[1], pt[2], pt[3]
        elif len(pt) == 3:
            x_, y_, lbl = pt; col = _GALLEDIA_HEX["RED"] if i == 0 else _GALLEDIA_HEX["G2"]
        else:
            x_, y_ = pt; lbl = ""; col = _GALLEDIA_HEX["RED"] if i == 0 else _GALLEDIA_HEX["G2"]
        ax.scatter(x_, y_, s=point_size, color=col, edgecolor="white", linewidth=1.5, zorder=3)
        if labels and lbl:
            ax.annotate(lbl, (x_, y_), xytext=(8, 8), textcoords="offset points",
                        fontsize=11, color="#000000")

    if x_label: ax.set_xlabel(x_label)
    if y_label: ax.set_ylabel(y_label)
    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_chart_scatter_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

# ── Sprint 4: Strategische Frameworks Teil 1 (v0.1.0-alpha.5 — IMPLEMENTED) ──

def _normalize_force(f):
    """Akzeptiert dict('title','items') oder tuple(title, items) oder list. Returns tuple."""
    if isinstance(f, dict):
        return f.get("title", ""), f.get("items", [])
    if isinstance(f, (tuple, list)) and len(f) == 2:
        return f[0], list(f[1]) if isinstance(f[1], (list, tuple)) else []
    return "", []

def swot(prs, strengths, weaknesses, opportunities, threats, kicker="", headline="", folio="", source=""):
    """SWOT-Analyse als 4-Quadranten-Layout.

    Args:
        strengths, weaknesses, opportunities, threats: list[str] — Bullet-Items pro Quadrant (max 5)

    Layout (klassisch MBB):
        ┌─────────────────┬─────────────────┐
        │ S  Strengths    │ W  Weaknesses   │  ← intern
        │  (TURK/Türkis)  │  (G3/Grau hell) │
        ├─────────────────┼─────────────────┤
        │ O Opportunities │ T  Threats      │  ← extern
        │  (RED/Galledia) │  (G1/Grau dunkel)│
        └─────────────────┴─────────────────┘
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    grid_x = _CHART_LEFT
    grid_y = _CHART_TOP
    grid_w = _CHART_WIDTH
    grid_h = _CHART_HEIGHT  # 6.8" ab v0.1.2 — Platz für 5+ Bullets pro Quadrant
    gap = Inches(0.18)
    cell_w = (grid_w - gap) // 2
    cell_h = (grid_h - gap) // 2

    cells = [
        ("S", "Strengths",     strengths,     TURK, WHITE, grid_x,                grid_y),
        ("W", "Weaknesses",    weaknesses,    G3,   BLACK, grid_x + cell_w + gap, grid_y),
        ("O", "Opportunities", opportunities, RED,  WHITE, grid_x,                grid_y + cell_h + gap),
        ("T", "Threats",       threats,       G1,   WHITE, grid_x + cell_w + gap, grid_y + cell_h + gap),
    ]

    for code, label, items, bg_color, text_color, x, y in cells:
        card(s, int(x), int(y), int(cell_w), int(cell_h), bg_color, corner=0.03)
        # Code-Letter (gross, links oben)
        code_tf = tb(s, int(x + Inches(0.25)), int(y + Inches(0.10)), Inches(1.0), Inches(0.85))
        run(code_tf.paragraphs[0], code, 44, SB, text_color, space_after=0)
        # Label rechts neben Code
        label_tf = tb(s, int(x + Inches(1.15)), int(y + Inches(0.30)), int(cell_w - Inches(1.3)), Inches(0.45))
        run(label_tf.paragraphs[0], label, 16, SB, text_color, space_after=0)
        # Bullet-Liste darunter
        body_tf = tb(s, int(x + Inches(0.30)), int(y + Inches(1.05)),
                     int(cell_w - Inches(0.6)), int(cell_h - Inches(1.2)))
        first = True
        for item in (items or [])[:5]:
            p = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
            run(p, f"•  {item}", 13, BODY, text_color, space_after=6)
            first = False

    _set_footer(s, prs, folio, source); return s

def porter_5f(prs, rivalry, new_entrants, substitutes, buyers, suppliers,
              kicker="", headline="", folio="", source=""):
    """Porter's Five Forces — zentrale Rivalry-Box mit 4 umgebenden Kräften.

    Args:
        rivalry: dict{'title','items'} oder tuple(title, items) — Branchen-Rivalität (Zentrum)
        new_entrants: dito — Bedrohung neuer Anbieter (oben)
        substitutes: dito — Bedrohung Substitute (unten)
        buyers: dito — Verhandlungsmacht Kunden (rechts)
        suppliers: dito — Verhandlungsmacht Lieferanten (links)

    Layout: Kreuz-Form, mit Pfeilen die alle zur zentralen Rivalry-Box zeigen.
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    cx_center = _CHART_LEFT + _CHART_WIDTH // 2
    cy_center = _CHART_TOP + _CHART_HEIGHT // 2

    # v0.1.3: Boxes nutzen jetzt die volle Inhaltszone (18.03" × 6.26")
    # Vertikal: outer(1.6) + spacing_v(0.4) + center(2.0) + spacing_v(0.4) + outer(1.6) = 6.0"
    # Horizontal: outer(6.0) + spacing_h(0.5) + center(5.0) + spacing_h(0.5) + outer(6.0) = 18.0"
    center_w = Inches(5.0)
    center_h = Inches(2.0)
    outer_w = Inches(6.0)
    outer_h = Inches(1.6)
    spacing = Inches(0.4)     # vertikal
    spacing_h = Inches(0.5)   # horizontal

    # Zentrum: Rivalität
    r_title, r_items = _normalize_force(rivalry)
    cx0 = int(cx_center - center_w // 2)
    cy0 = int(cy_center - center_h // 2)
    card(s, cx0, cy0, int(center_w), int(center_h), RED, corner=0.06)
    rtf = tb(s, cx0 + int(Inches(0.3)), cy0 + int(Inches(0.12)),
             int(center_w - Inches(0.6)), int(center_h - Inches(0.24)))
    run(rtf.paragraphs[0], "Rivalität in der Branche", 18, SB, WHITE, align=PP_ALIGN.CENTER, space_after=3)
    if r_title:
        run(rtf.add_paragraph(), r_title, 13, BODY, WHITE, align=PP_ALIGN.CENTER, space_after=5)
    for item in (r_items or [])[:3]:
        run(rtf.add_paragraph(), f"•  {item}", 12, BODY, WHITE, align=PP_ALIGN.LEFT, space_after=2)

    # 4 Outer Forces — (default_title, force_data, dx, dy, arrow_direction)
    outer_specs = [
        ("Neue Anbieter",   new_entrants, 0,
         -(center_h // 2 + outer_h // 2 + spacing), "down"),
        ("Substitute",      substitutes,  0,
          (center_h // 2 + outer_h // 2 + spacing), "up"),
        ("Verhandlungsmacht Kunden",  buyers,
          (center_w // 2 + outer_w // 2 + spacing_h), 0, "left"),
        ("Verhandlungsmacht Lieferanten", suppliers,
         -(center_w // 2 + outer_w // 2 + spacing_h), 0, "right"),
    ]

    from pptx.enum.shapes import MSO_CONNECTOR
    from lxml import etree
    for default_title, force_data, dx, dy, arrow_dir in outer_specs:
        f_title, f_items = _normalize_force(force_data)
        if not f_title: f_title = default_title
        x = int(cx_center + dx - outer_w // 2)
        y = int(cy_center + dy - outer_h // 2)
        card(s, x, y, int(outer_w), int(outer_h), GL, corner=0.05)
        otf = tb(s, x + int(Inches(0.25)), y + int(Inches(0.15)),
                 int(outer_w - Inches(0.5)), int(outer_h - Inches(0.3)))
        run(otf.paragraphs[0], f_title, 15, SB, BLACK, space_after=4)
        for item in (f_items or [])[:3]:
            run(otf.add_paragraph(), f"•  {item}", 11, BODY, G1, space_after=2)

        # Pfeil zum Zentrum
        if arrow_dir == "down":
            start = (cx_center, y + int(outer_h))
            end   = (cx_center, cy_center - center_h // 2)
        elif arrow_dir == "up":
            start = (cx_center, y)
            end   = (cx_center, cy_center + center_h // 2)
        elif arrow_dir == "left":
            start = (x, cy_center)
            end   = (cx_center + center_w // 2, cy_center)
        else:  # right
            start = (x + int(outer_w), cy_center)
            end   = (cx_center - center_w // 2, cy_center)

        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       int(start[0]), int(start[1]), int(end[0]), int(end[1]))
        conn.line.color.rgb = G2; conn.line.width = Pt(2)
        ln = conn.line._get_or_add_ln()
        tailEnd = ln.find(qn('a:tailEnd'))
        if tailEnd is None:
            tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle'); tailEnd.set('w', 'med'); tailEnd.set('len', 'med')

    _set_footer(s, prs, folio, source); return s

def value_chain(prs, primary_activities, support_activities,
                kicker="", headline="", folio="", source=""):
    """Porter's Value Chain — Stützfunktionen (oben gestapelt) + Primärtätigkeiten (unten als Pfeilkette).

    Args:
        primary_activities: list[tuple[title, items]] — klassisch 5
            (Inbound Logistics, Operations, Outbound Logistics, Marketing & Sales, Service)
        support_activities: list[tuple[title, items]] — klassisch 4
            (Firm Infrastructure, HRM, Technology Development, Procurement)

    Layout:
        ┌── Firm Infrastructure ─────────────────────┐
        ├── HRM ────────────────────────────────────┤   Support
        ├── Technology Development ─────────────────┤   (Top)
        ├── Procurement ────────────────────────────┤
        └────────────────────────────────────────────┘
        ▶ Inbound ▶ Ops ▶ Outbound ▶ M&S ▶ Service ▶ MARGIN
                              Primary (Bottom)
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    chart_x = _CHART_LEFT
    chart_y = _CHART_TOP
    chart_w = _CHART_WIDTH - Inches(1.4)  # rechts Platz für "Margin"
    chart_h = _CHART_HEIGHT

    # Support: oben, 4 horizontale Balken
    n_sup = len(support_activities)
    sup_total_h = int(chart_h * 0.45)
    sup_gap = int(Inches(0.06))
    bar_h = (sup_total_h - (n_sup - 1) * sup_gap) // max(1, n_sup)

    for i, sup in enumerate(support_activities):
        title, items = _normalize_force(sup) if isinstance(sup, dict) else (sup if isinstance(sup, tuple) and len(sup) == 2 else ("", []))
        if isinstance(sup, tuple) and len(sup) == 2:
            title, items = sup[0], list(sup[1])
        y = chart_y + i * (bar_h + sup_gap)
        card(s, int(chart_x), int(y), int(chart_w), int(bar_h), G2, corner=0.02)
        # Title links
        ttf = tb(s, int(chart_x + Inches(0.25)), int(y), Inches(3.2), bar_h, MSO_ANCHOR.MIDDLE)
        run(ttf.paragraphs[0], title, 13, SB, WHITE, space_after=0)
        # Items rechts (kompakt, kommagetrennt)
        if items:
            items_text = " · ".join(items[:4])
            itf = tb(s, int(chart_x + Inches(3.5)), int(y), int(chart_w - Inches(3.7)), bar_h, MSO_ANCHOR.MIDDLE)
            run(itf.paragraphs[0], items_text, 11, BODY, WHITE, space_after=0)

    # Primary: unten, Pfeilkette aus Rechtecken + Trennpfeilen
    prim_y = chart_y + sup_total_h + int(Inches(0.30))
    prim_h = int(chart_h * 0.45)
    n_prim = len(primary_activities)
    arrow_w = int(Inches(0.35))
    seg_w = (chart_w - (n_prim - 1) * arrow_w) // max(1, n_prim)

    for i, prim in enumerate(primary_activities):
        if isinstance(prim, tuple) and len(prim) == 2:
            title, items = prim[0], list(prim[1])
        else:
            title, items = _normalize_force(prim)
        x = chart_x + i * (seg_w + arrow_w)
        is_last = (i == n_prim - 1)
        fill = RED if is_last else GL
        text_color = WHITE if is_last else BLACK
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(prim_y), int(seg_w), int(prim_h))
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = WHITE; rect.line.width = Pt(1.5)
        rect.shadow.inherit = False

        # Title oben, Items kompakt darunter
        ttf = tb(s, int(x + Inches(0.15)), int(prim_y + Inches(0.15)),
                 int(seg_w - Inches(0.3)), Inches(0.5))
        run(ttf.paragraphs[0], title, 12, SB, text_color, align=PP_ALIGN.CENTER, space_after=4)
        if items:
            itf = tb(s, int(x + Inches(0.15)), int(prim_y + Inches(0.65)),
                     int(seg_w - Inches(0.3)), int(prim_h - Inches(0.75)))
            first = True
            for item in items[:3]:
                p = itf.paragraphs[0] if first else itf.add_paragraph()
                run(p, item, 10, BODY, text_color, align=PP_ALIGN.CENTER, space_after=3)
                first = False

        # Pfeil-Spitze zum nächsten Segment (Right-Arrow zwischen Rechtecken)
        if not is_last:
            arrow_x = x + seg_w
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       int(arrow_x), int(prim_y + prim_h // 2 - Inches(0.20)),
                                       arrow_w, int(Inches(0.40)))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = G3
            arrow.line.fill.background(); arrow.shadow.inherit = False

    # Margin-Label rechts
    margin_x = chart_x + chart_w + int(Inches(0.15))
    margin_tf = tb(s, int(margin_x), int(prim_y), Inches(1.2), int(prim_h), MSO_ANCHOR.MIDDLE)
    run(margin_tf.paragraphs[0], "Margin", 16, SB, RED, align=PP_ALIGN.CENTER, space_after=0)

    _set_footer(s, prs, folio, source); return s

def bcg_matrix(prs, products, kicker="", headline="", folio="", source=""):
    """BCG-Portfolio-Matrix — Marktanteil (x, invertiert) × Marktwachstum (y).

    Args:
        products: list[dict] mit:
            "name": str (Produkt/SBU-Name)
            "share": float in [0,1] — relativer Marktanteil (1=Marktführer, 0=Nachzügler)
            "growth": float in [0,1] — Marktwachstum (1=hoch, 0=niedrig)
            "revenue": float (optional) — Umsatz, bestimmt Bubble-Grösse

    Quadranten:
        Stars (★)         Top-Left  | Question Marks (?) Top-Right
        Cash Cows ($)     Bottom-L  | Dogs (✗)           Bottom-R

    Klassische BCG-Konvention: X-Achse invertiert (hoher Marktanteil = LINKS).
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    qx = _CHART_LEFT + Inches(1.5)
    qy = _CHART_TOP
    qw = _CHART_WIDTH - Inches(2.4)
    qh = _CHART_HEIGHT - Inches(0.5)
    half_w = qw // 2
    half_h = qh // 2

    # 4 Quadranten mit hellen Hintergründen + Icon-Labels
    quads = [
        ("Stars",          "★", qx,            qy,            half_w, half_h, RGBColor(0xCC, 0xEE, 0xEC)),
        ("Question Marks", "?", qx + half_w,   qy,            half_w, half_h, RGBColor(0xFD, 0xE7, 0xED)),
        ("Cash Cows",      "$", qx,            qy + half_h,   half_w, half_h, GL),
        ("Dogs",           "✗", qx + half_w,   qy + half_h,   half_w, half_h, GL),
    ]
    for label, icon, x_, y_, w_, h_, fill in quads:
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x_), int(y_), int(w_), int(h_))
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = G4; rect.line.width = Pt(0.75)
        rect.shadow.inherit = False
        ltf = tb(s, int(x_ + Inches(0.20)), int(y_ + Inches(0.12)), Inches(3.5), Inches(0.45))
        run(ltf.paragraphs[0], f"{icon}  {label}", 14, SB, BLACK, space_after=0)

    # Achsenbeschriftungen (X invertiert: hoch links, niedrig rechts!)
    x_lbl_l = tb(s, int(qx), int(qy + qh + Inches(0.10)), Inches(3.5), Inches(0.3))
    run(x_lbl_l.paragraphs[0], "Hoher Marktanteil", 11, BODY, G2, align=PP_ALIGN.LEFT, space_after=0)
    x_lbl_r = tb(s, int(qx + qw - Inches(3.5)), int(qy + qh + Inches(0.10)), Inches(3.5), Inches(0.3))
    run(x_lbl_r.paragraphs[0], "Niedriger Marktanteil", 11, BODY, G2, align=PP_ALIGN.RIGHT, space_after=0)
    y_lbl_t = tb(s, int(qx - Inches(1.4)), int(qy + Inches(0.10)), Inches(1.3), Inches(0.3))
    run(y_lbl_t.paragraphs[0], "Hohes Wachstum", 11, BODY, G2, align=PP_ALIGN.RIGHT, space_after=0)
    y_lbl_b = tb(s, int(qx - Inches(1.4)), int(qy + qh - Inches(0.40)), Inches(1.3), Inches(0.3))
    run(y_lbl_b.paragraphs[0], "Niedriges Wachstum", 11, BODY, G2, align=PP_ALIGN.RIGHT, space_after=0)

    # Bubbles (Galledia-Rot, Grösse via Umsatz)
    max_rev = max((p.get("revenue", 1) for p in products), default=1) or 1
    for prod in products:
        share = max(0.0, min(1.0, prod.get("share", 0.5)))
        growth = max(0.0, min(1.0, prod.get("growth", 0.5)))
        revenue = prod.get("revenue", 1)
        name = prod.get("name", "?")
        bubble_d = int(Inches(0.40 + 0.80 * (revenue / max_rev)))

        # Klassische BCG-Konvention: share=1 (hoch) → ix=0 (links), share=0 (niedrig) → ix=1 (rechts)
        ix = 1 - share
        iy = growth  # growth=1 → oben

        px = int(qx + ix * qw) - bubble_d // 2
        py = int(qy + (1 - iy) * qh) - bubble_d // 2

        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, px, py, bubble_d, bubble_d)
        circle.fill.solid(); circle.fill.fore_color.rgb = RED
        circle.line.color.rgb = WHITE; circle.line.width = Pt(2)
        circle.shadow.inherit = False

        # Label
        ltf = tb(s, px + bubble_d + int(Inches(0.05)), py - int(Inches(0.02)),
                 int(Inches(2.5)), int(Inches(0.4)))
        run(ltf.paragraphs[0], name, 11, SB, BLACK, align=PP_ALIGN.LEFT, space_after=0)

    _set_footer(s, prs, folio, source); return s

# ── Sprint 5: Strategische Frameworks Teil 2 (v0.1.0-alpha.6 — IMPLEMENTED) ──

def seven_s(prs, items, kicker="", headline="", folio="", source=""):
    """McKinsey 7S Framework — Shared Values im Zentrum, 6 S's drumherum.

    Args:
        items: dict mit 7 Keys (jeweils dict mit 'title' und optional 'items'):
            "shared_values": (Zentrum, RED, grösster)
            "strategy", "structure", "systems"      → "Hard S's" (oben)
            "style", "staff", "skills"               → "Soft S's" (unten)

    Beispiel:
        seven_s(prs, items={
            "shared_values": {"title": "Innovation", "items": ["Pionier", "Mut"]},
            "strategy":      {"title": "Plattform-Bündel", "items": ["Print + Digital"]},
            "structure":     {"title": "5 OE", "items": ["Holding-Modell"]},
            "systems":       {"title": "Twenty CRM", "items": ["n8n", "AI Hub"]},
            "skills":        {"title": "Mediaberatung", "items": ["Editorial"]},
            "style":         {"title": "Pragmatisch", "items": []},
            "staff":         {"title": "ca. 200 MA", "items": []},
        }, kicker="Organisations-Diagnose", headline="McKinsey 7S Galledia")
    """
    import math
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    cx = _CHART_LEFT + _CHART_WIDTH // 2
    cy = _CHART_TOP + _CHART_HEIGHT // 2 - int(Inches(0.1))
    radius = int(Inches(2.2))
    center_d = int(Inches(2.4))
    outer_d  = int(Inches(2.0))

    # Outer-S-Positionen (6 Stück, Uhrzeiger-Stunden 1,3,5,7,9,11)
    # Hard S's oben, Soft S's unten
    layout = [
        # (key, label, angle_deg from horizontal, color)
        ("strategy",   "Strategy",   60,  TURK),   # oben rechts
        ("skills",     "Skills",      0,  TURK),   # rechts mitte
        ("style",      "Style",     -60,  G2),     # unten rechts
        ("staff",      "Staff",    -120,  G2),     # unten links
        ("systems",    "Systems",  -180,  TURK),   # links mitte
        ("structure",  "Structure", 120,  TURK),   # oben links
    ]

    # 1. Verbindungslinien Zentrum → Outer (zuerst, damit sie unter den Kreisen liegen)
    from pptx.enum.shapes import MSO_CONNECTOR
    for key, label, angle, color in layout:
        rad = math.radians(angle)
        ox = cx + int(radius * math.cos(rad))
        oy = cy - int(radius * math.sin(rad))
        line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(cx), int(cy), int(ox), int(oy))
        line.line.color.rgb = G3; line.line.width = Pt(1.2)

    # 2. Zentrum: Shared Values
    sv = items.get("shared_values", {})
    sv_title = sv.get("title", "Shared Values") if isinstance(sv, dict) else (sv[0] if isinstance(sv, tuple) else "Shared Values")
    sv_items = sv.get("items", []) if isinstance(sv, dict) else (list(sv[1]) if isinstance(sv, tuple) and len(sv) > 1 else [])
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 int(cx - center_d // 2), int(cy - center_d // 2),
                                 center_d, center_d)
    circle.fill.solid(); circle.fill.fore_color.rgb = RED
    circle.line.color.rgb = WHITE; circle.line.width = Pt(3)
    circle.shadow.inherit = False
    sv_tf = tb(s, int(cx - center_d // 2 + Inches(0.1)), int(cy - Inches(0.45)),
               int(center_d - Inches(0.2)), Inches(0.9), MSO_ANCHOR.MIDDLE)
    run(sv_tf.paragraphs[0], "Shared Values", 10, BODY, WHITE, align=PP_ALIGN.CENTER, space_after=2)
    run(sv_tf.add_paragraph(), sv_title, 14, SB, WHITE, align=PP_ALIGN.CENTER, space_after=0)

    # 3. Outer-S-Kreise mit Title
    for key, label, angle, color in layout:
        rad = math.radians(angle)
        ox = cx + int(radius * math.cos(rad))
        oy = cy - int(radius * math.sin(rad))
        # Daten holen
        d = items.get(key, {})
        if isinstance(d, dict):
            d_title = d.get("title", "")
            d_items = d.get("items", [])
        elif isinstance(d, tuple) and len(d) >= 1:
            d_title = d[0]; d_items = list(d[1]) if len(d) > 1 else []
        else:
            d_title = ""; d_items = []
        # Kreis
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     int(ox - outer_d // 2), int(oy - outer_d // 2),
                                     outer_d, outer_d)
        circle.fill.solid(); circle.fill.fore_color.rgb = color
        circle.line.color.rgb = WHITE; circle.line.width = Pt(2.5)
        circle.shadow.inherit = False
        # Text im Kreis (Label oben, Title gross darunter)
        ttf = tb(s, int(ox - outer_d // 2 + Inches(0.1)), int(oy - Inches(0.5)),
                 int(outer_d - Inches(0.2)), Inches(1.0), MSO_ANCHOR.MIDDLE)
        run(ttf.paragraphs[0], label, 10, BODY, WHITE, align=PP_ALIGN.CENTER, space_after=2)
        if d_title:
            run(ttf.add_paragraph(), d_title, 12, SB, WHITE, align=PP_ALIGN.CENTER, space_after=0)

    _set_footer(s, prs, folio, source); return s

def ansoff_matrix(prs, strategies, kicker="", headline="", folio="", source=""):
    """Ansoff-Wachstumsmatrix: Produkte (y) × Märkte (x), 4 Wachstumsstrategien.

    Args:
        strategies: dict mit 4 Keys (jeweils list[str] mit Bullet-Items):
            "market_penetration"    — Existing × Existing (TURK, geringstes Risiko)
            "product_development"   — Existing market × New product (G1)
            "market_development"    — New market × Existing product (G2)
            "diversification"       — New × New (RED, höchstes Risiko)

    Layout:
        ┌──────────────────┬──────────────────┐
        │ Product          │ Diversification  │  ← New Products
        │ Development (G1) │ (RED)            │
        ├──────────────────┼──────────────────┤
        │ Market           │ Market           │  ← Existing Products
        │ Penetration(TURK)│ Development (G2) │
        └──────────────────┴──────────────────┘
          Existing Markets    New Markets
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    grid_x = _CHART_LEFT + Inches(1.4)
    grid_y = _CHART_TOP
    grid_w = _CHART_WIDTH - Inches(2.0)
    grid_h = _CHART_HEIGHT - Inches(0.5)
    gap = Inches(0.18)
    cell_w = (grid_w - gap) // 2
    cell_h = (grid_h - gap) // 2

    cells = [
        # (key, title, risk_label, bg_color, text_color, x, y)
        ("product_development",   "Product Development",  "Mittel",   G1,   WHITE,
            grid_x,                grid_y),
        ("diversification",       "Diversification",      "Hoch",     RED,  WHITE,
            grid_x + cell_w + gap, grid_y),
        ("market_penetration",    "Market Penetration",   "Niedrig",  TURK, WHITE,
            grid_x,                grid_y + cell_h + gap),
        ("market_development",    "Market Development",   "Mittel",   G2,   WHITE,
            grid_x + cell_w + gap, grid_y + cell_h + gap),
    ]

    for key, title, risk, bg, text_color, x, y in cells:
        card(s, int(x), int(y), int(cell_w), int(cell_h), bg, corner=0.03)
        # Header: Strategie-Name + Risiko-Tag
        ttf = tb(s, int(x + Inches(0.30)), int(y + Inches(0.20)),
                 int(cell_w - Inches(0.6)), Inches(0.5))
        run(ttf.paragraphs[0], title, 16, SB, text_color, space_after=2)
        risk_tf = tb(s, int(x + Inches(0.30)), int(y + Inches(0.80)),
                     int(cell_w - Inches(0.6)), Inches(0.30))
        run(risk_tf.paragraphs[0], f"Risiko: {risk}", 10, BODY, text_color, italic=True, space_after=0)
        # Bullet-Liste
        bullets = strategies.get(key, []) or []
        if bullets:
            body_tf = tb(s, int(x + Inches(0.30)), int(y + Inches(1.30)),
                         int(cell_w - Inches(0.6)), int(cell_h - Inches(1.5)))
            first = True
            for item in bullets[:4]:
                p = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
                run(p, f"•  {item}", 12, BODY, text_color, space_after=5)
                first = False

    # Achsen-Labels
    # X-Achse (Märkte)
    x_l = tb(s, int(grid_x), int(grid_y + grid_h + Inches(0.10)), int(cell_w), Inches(0.30))
    run(x_l.paragraphs[0], "Bestehende Märkte", 12, SB, G2, align=PP_ALIGN.CENTER, space_after=0)
    x_r = tb(s, int(grid_x + cell_w + gap), int(grid_y + grid_h + Inches(0.10)), int(cell_w), Inches(0.30))
    run(x_r.paragraphs[0], "Neue Märkte", 12, SB, G2, align=PP_ALIGN.CENTER, space_after=0)
    # Y-Achse (Produkte) — links neben Grid
    y_t = tb(s, int(grid_x - Inches(1.35)), int(grid_y + cell_h // 2 - Inches(0.20)),
             Inches(1.25), Inches(0.40))
    run(y_t.paragraphs[0], "Neue\nProdukte", 11, SB, G2, align=PP_ALIGN.RIGHT, space_after=0)
    y_b = tb(s, int(grid_x - Inches(1.35)), int(grid_y + cell_h + gap + cell_h // 2 - Inches(0.20)),
             Inches(1.25), Inches(0.40))
    run(y_b.paragraphs[0], "Bestehende\nProdukte", 11, SB, G2, align=PP_ALIGN.RIGHT, space_after=0)

    _set_footer(s, prs, folio, source); return s

def maturity_model(prs, levels, current_level=None, target_level=None,
                   kicker="", headline="", folio="", source=""):
    """Maturity-Model als aufsteigende Stufentreppe.

    Args:
        levels: list[dict] mit 'name' und optional 'description'
                Beispiel: [{"name": "Initial"}, {"name": "Repeatable"}, ...]
                Typisch 5 Levels (CMMI-Standard) oder 3-7.
        current_level: int (1-based) — aktueller Reifegrad (RED-Marker)
        target_level: int (1-based) — Ziel-Reifegrad (TURK-Marker)

    Layout: jede Stufe ist ein Rechteck, das nach rechts höher wird.
            Current = roter Pfeil von unten, Target = türkiser Stern.

    Beispiel:
        maturity_model(prs,
            levels=[
                {"name": "Initial",     "description": "Ad-hoc, reaktiv"},
                {"name": "Repeatable",  "description": "Erste Standards"},
                {"name": "Defined",     "description": "Dokumentierte Prozesse"},
                {"name": "Managed",     "description": "Messbar gesteuert"},
                {"name": "Optimized",   "description": "Kontinuierl. Verbesserung"},
            ],
            current_level=2, target_level=4,
            kicker="CRM-Reife", headline="Galledia Q1 → Q4 2026")
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    n = len(levels)
    chart_x = _CHART_LEFT
    chart_y = _CHART_TOP + Inches(0.4)  # Platz oben für Marker
    chart_w = _CHART_WIDTH
    chart_h = _CHART_HEIGHT - Inches(1.2)  # Platz unten für Labels

    step_w = chart_w // n
    # Step-Höhen: linear aufsteigend, 25% bis 100% von chart_h
    min_h_frac, max_h_frac = 0.25, 1.0
    step_heights = []
    for i in range(n):
        frac = min_h_frac + (max_h_frac - min_h_frac) * (i / max(1, n - 1))
        step_heights.append(int(chart_h * frac))

    # Stufen zeichnen (links unten = niedrig, rechts unten = hoch)
    base_y = chart_y + chart_h
    for i, lvl in enumerate(levels):
        x = chart_x + i * step_w
        h = step_heights[i]
        y = base_y - h
        # Farbe: Standard G3, current = RED, target = TURK
        is_current = (current_level is not None and current_level == i + 1)
        is_target  = (target_level  is not None and target_level  == i + 1)
        if is_current:
            fill = RED
        elif is_target:
            fill = TURK
        else:
            fill = G3
        text_color = WHITE
        rect = card(s, int(x), int(y), int(step_w - Inches(0.05)), int(h), fill, corner=0.02)

        # Number-Badge oben links
        num_tf = tb(s, int(x + Inches(0.15)), int(y + Inches(0.10)), Inches(0.7), Inches(0.5))
        run(num_tf.paragraphs[0], f"{i+1}", 28, SB, text_color, space_after=0)
        # Name darunter
        name = lvl.get("name", "") if isinstance(lvl, dict) else str(lvl)
        name_tf = tb(s, int(x + Inches(0.15)), int(y + Inches(0.75)),
                     int(step_w - Inches(0.3)), Inches(0.5))
        run(name_tf.paragraphs[0], name, 13, SB, text_color, space_after=2)
        # Description (falls vorhanden, nur wenn Stufe gross genug)
        if isinstance(lvl, dict) and lvl.get("description") and h > int(Inches(1.5)):
            desc_tf = tb(s, int(x + Inches(0.15)), int(y + Inches(1.30)),
                         int(step_w - Inches(0.3)), int(h - Inches(1.4)))
            run(desc_tf.paragraphs[0], lvl["description"], 10, BODY, text_color, space_after=0)

    # Marker: Current (rote "HEUTE"-Box oberhalb)
    if current_level is not None:
        x = chart_x + (current_level - 1) * step_w
        h = step_heights[current_level - 1]
        y = base_y - h
        marker_tf = tb(s, int(x), int(y - Inches(0.45)), int(step_w - Inches(0.05)), Inches(0.35))
        run(marker_tf.paragraphs[0], "▼ HEUTE", 12, SB, RED, align=PP_ALIGN.CENTER, space_after=0)
    if target_level is not None and target_level != current_level:
        x = chart_x + (target_level - 1) * step_w
        h = step_heights[target_level - 1]
        y = base_y - h
        marker_tf = tb(s, int(x), int(y - Inches(0.45)), int(step_w - Inches(0.05)), Inches(0.35))
        run(marker_tf.paragraphs[0], "▼ ZIEL", 12, SB, TURK, align=PP_ALIGN.CENTER, space_after=0)

    _set_footer(s, prs, folio, source); return s

def three_horizons(prs, h1_now, h2_emerging, h3_future,
                   kicker="", headline="", folio="", source=""):
    """McKinsey 3-Horizons Framework — überlappende Wachstumskurven über Zeit.

    Args:
        h1_now:      dict{'title','items'} — Kerngeschäft, jetzt-1J (sinkende Kurve)
        h2_emerging: dict{'title','items'} — Aufkommende Geschäfte, 2-4J
        h3_future:   dict{'title','items'} — Zukunftsoptionen, 4+J

    Layout: Zeit-Achse (x) horizontal, Performance (y) vertikal.
            3 gefärbte Bereiche, jeweils mit Text-Box rechts neben dem Peak.
            H1 fällt ab, H2 ist mittlerer Bogen, H3 steigt.
    """
    import math
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    chart_x = _CHART_LEFT + Inches(1.2)
    chart_y = _CHART_TOP
    chart_w = _CHART_WIDTH - Inches(1.5)
    chart_h = _CHART_HEIGHT - Inches(0.4)

    # Y-Achse (Performance/Wert) — Pfeil oben
    from pptx.enum.shapes import MSO_CONNECTOR
    from lxml import etree
    y_axis = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     int(chart_x), int(chart_y + chart_h),
                                     int(chart_x), int(chart_y))
    y_axis.line.color.rgb = G2; y_axis.line.width = Pt(1.5)
    ln = y_axis.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')

    # X-Achse (Zeit) — Pfeil rechts
    x_axis = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     int(chart_x), int(chart_y + chart_h),
                                     int(chart_x + chart_w), int(chart_y + chart_h))
    x_axis.line.color.rgb = G2; x_axis.line.width = Pt(1.5)
    ln2 = x_axis.line._get_or_add_ln()
    tail2 = etree.SubElement(ln2, qn('a:tailEnd'))
    tail2.set('type', 'triangle'); tail2.set('w', 'med'); tail2.set('len', 'med')

    # Achsen-Labels
    y_lbl = tb(s, int(chart_x - Inches(1.1)), int(chart_y), Inches(1.0), Inches(0.4))
    run(y_lbl.paragraphs[0], "Wert", 11, SB, G2, align=PP_ALIGN.RIGHT, space_after=0)
    x_lbl = tb(s, int(chart_x + chart_w - Inches(0.6)), int(chart_y + chart_h + Inches(0.10)),
               Inches(0.6), Inches(0.30))
    run(x_lbl.paragraphs[0], "Zeit", 11, SB, G2, align=PP_ALIGN.RIGHT, space_after=0)

    # 3 Horizons als Freeform-Bögen — vereinfacht als gefärbte Ovale (Wellen)
    # H1: sinkt von hoch nach niedrig (linke Seite)
    # H2: mittlere Welle (Mitte)
    # H3: steigt von niedrig nach hoch (rechte Seite)
    # Wir nutzen 3 OVAL-Shapes mit Position/Grösse für den Effekt überlappender Wellen.

    horizons = [
        # (data, x_center_frac, y_center_frac, w_frac, h_frac, color)
        (h1_now,      0.20, 0.40, 0.45, 0.45, G1),    # links, mittlere Höhe, sinkende Aufmerksamkeit
        (h2_emerging, 0.50, 0.55, 0.45, 0.55, TURK),  # Mitte
        (h3_future,   0.80, 0.30, 0.45, 0.65, RED),   # rechts, am höchsten
    ]
    for i, (data, cx_frac, cy_frac, w_frac, h_frac, color) in enumerate(horizons):
        w = int(chart_w * w_frac)
        h = int(chart_h * h_frac)
        cx_h = chart_x + int(chart_w * cx_frac)
        cy_h = chart_y + int(chart_h * (1 - cy_frac))
        # Halbtransparente Wirkung simulieren via OVAL mit Linie statt Fill
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   cx_h - w // 2, cy_h - h // 2, w, h)
        oval.fill.solid(); oval.fill.fore_color.rgb = color
        oval.line.color.rgb = WHITE; oval.line.width = Pt(3)
        oval.shadow.inherit = False

        # Horizont-Label im Oval (mittig)
        title, items = _normalize_force(data)
        h_idx = i + 1
        label_tf = tb(s, cx_h - int(Inches(2.0)), cy_h - int(Inches(0.45)),
                      int(Inches(4.0)), int(Inches(0.9)), MSO_ANCHOR.MIDDLE)
        run(label_tf.paragraphs[0], f"Horizont {h_idx}", 12, BODY, WHITE,
            align=PP_ALIGN.CENTER, space_after=2)
        if title:
            run(label_tf.add_paragraph(), title, 16, SB, WHITE,
                align=PP_ALIGN.CENTER, space_after=0)

    # Zeit-Labels für Horizonte unter X-Achse
    time_labels = [("Jetzt → 1J", 0.20), ("2–4 Jahre", 0.50), ("4+ Jahre", 0.80)]
    for label, frac in time_labels:
        tlx = chart_x + int(chart_w * frac) - int(Inches(1.0))
        tl_tf = tb(s, tlx, int(chart_y + chart_h + Inches(0.40)),
                   int(Inches(2.0)), int(Inches(0.30)))
        run(tl_tf.paragraphs[0], label, 11, BODY, G1, align=PP_ALIGN.CENTER, space_after=0)

    # Inhalte unterhalb der Chart-Area (3 kleine Bullet-Listen)
    detail_y = chart_y + chart_h + int(Inches(0.9))
    detail_h = int(_CHART_HEIGHT) - (detail_y - int(_CHART_TOP)) - int(Inches(0.1))
    if detail_h > int(Inches(0.5)):
        col_w = chart_w // 3
        for i, (data, _, _, _, _, color) in enumerate(horizons):
            _, items = _normalize_force(data)
            if items:
                dx = chart_x + i * col_w
                d_tf = tb(s, int(dx + Inches(0.1)), detail_y,
                          int(col_w - Inches(0.2)), int(detail_h))
                first = True
                for item in items[:3]:
                    p = d_tf.paragraphs[0] if first else d_tf.add_paragraph()
                    run(p, f"•  {item}", 10, BODY, G1, space_after=2)
                    first = False

    _set_footer(s, prs, folio, source); return s

# ── Sprint 6: Storyline-Enforcement (v0.1.0-alpha.7 — IMPLEMENTED) ───────────

def exec_summary(prs, headline, key_points, kicker="Executive Summary", folio="", source=""):
    """Executive-Summary-Folie — Hauptaussage + 3-5 Stützpunkte mit Folien-Referenzen.

    Args:
        prs: Presentation
        headline: str — Kernaussage des Decks in 1 Satz (max 32 Zeichen wie üblich)
        key_points: list[dict] mit "text" und optional "slide_ref":
            [{"text": "Kosten -23%", "slide_ref": "S. 7"}, ...]
            ODER list[tuple[text, slide_ref]]
            ODER list[str] (ohne Refs)

    Layout: nummerierte Stützpunkte (1-5), pro Punkt grosse Zahl in RED,
            Text rechts daneben, Folien-Referenz dezent in G2.

    Beispiel:
        exec_summary(prs,
            headline="AI-Hub zahlt sich aus",
            key_points=[
                {"text": "3'900 h/Monat eingespart", "slide_ref": "S. 7"},
                {"text": "Amortisation in 18 Mt",    "slide_ref": "S. 11"},
                {"text": "5 von 5 Sparten profitieren","slide_ref": "S. 14"},
            ])
    """
    s = _blank(prs)
    _set_header(s, kicker, headline)

    # Normalisiere key_points
    normalized = []
    for kp in key_points:
        if isinstance(kp, dict):
            normalized.append((kp.get("text", ""), kp.get("slide_ref", "")))
        elif isinstance(kp, tuple) and len(kp) == 2:
            normalized.append((kp[0], kp[1]))
        elif isinstance(kp, tuple) and len(kp) == 1:
            normalized.append((kp[0], ""))
        else:
            normalized.append((str(kp), ""))

    n = len(normalized)
    if n == 0: return s

    item_h = Inches(0.95)
    gap = Inches(0.15)
    total_h = n * item_h + (n - 1) * gap
    start_y = _CHART_TOP + (_CHART_HEIGHT - total_h) // 2

    for i, (text, ref) in enumerate(normalized):
        y = start_y + i * (item_h + gap)
        # Nummer-Badge links (rot, Kreis)
        badge_d = int(Inches(0.85))
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    int(_CHART_LEFT), int(y),
                                    badge_d, badge_d)
        badge.fill.solid(); badge.fill.fore_color.rgb = RED
        badge.line.fill.background(); badge.shadow.inherit = False
        bf = tb(s, int(_CHART_LEFT), int(y), badge_d, badge_d, MSO_ANCHOR.MIDDLE)
        run(bf.paragraphs[0], str(i + 1), 36, SB, WHITE, align=PP_ALIGN.CENTER, space_after=0)

        # Text rechts neben Badge
        text_x = _CHART_LEFT + badge_d + int(Inches(0.35))
        text_w = _CHART_WIDTH - badge_d - int(Inches(0.35))
        if ref:
            text_w -= int(Inches(1.5))  # Platz für Ref
        ttf = tb(s, int(text_x), int(y + Inches(0.05)), int(text_w), int(item_h - Inches(0.1)), MSO_ANCHOR.MIDDLE)
        run(ttf.paragraphs[0], text, 22, SB, BLACK, space_after=0)

        # Slide-Ref rechts (dezent)
        if ref:
            rf = tb(s, int(_CHART_LEFT + _CHART_WIDTH - Inches(1.4)), int(y + Inches(0.05)),
                    Inches(1.3), int(item_h - Inches(0.1)), MSO_ANCHOR.MIDDLE)
            run(rf.paragraphs[0], f"→ {ref}", 13, BODY, G2, align=PP_ALIGN.RIGHT, italic=True, space_after=0)

    _set_footer(s, prs, folio, source); return s

def add_governing_thought(slide, text, y_pos=None, color=None, text_color=None, size_pt=13):
    """Governing-Thought-Banner — 1-Satz-Argument-Zusammenfassung als horizontaler Bar.

    Args:
        slide: Slide-Objekt (existierende Folie!)
        text: 1-Satz-Aussage der Folie (z.B. "Die Print-Schwäche wird durch Digital überkompensiert")
        y_pos: y-Koordinate in EMU (Default Inches(2.40), zwischen Headline und Content)
        color: Hintergrund-Farbe (Default TURK)
        text_color: Schriftfarbe (Default WHITE)
        size_pt: Schriftgrösse (Default 13)

    Sitzt zwischen Headline (y≈1.06") und Content-Bereich (y≈3.30"), unterstreicht
    die Folie-Argumentation als „governing thought" im MBB-Sinne.
    """
    if y_pos is None: y_pos = Inches(2.40)
    if color is None: color = TURK
    if text_color is None: text_color = WHITE

    card(slide, int(GX), int(y_pos), int(GW), int(Inches(0.55)), color, corner=0.03)
    tf = tb(slide, int(GX + Inches(0.30)), int(y_pos),
            int(GW - Inches(0.6)), int(Inches(0.55)), MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], text, size_pt, BODY, text_color,
        align=PP_ALIGN.LEFT, italic=True, space_after=0)
    return slide

def add_section_tracker(slide, sections, current_idx, show_breadcrumb=False, y_pos=None, override_kicker=True):
    """Section-Tracker — setzt den Kapiteltitel auf die aktive Section.

    Ab v0.1.2: Standardverhalten ist NUR Kapiteltitel-Override (kein Breadcrumb).
    Die früheren Bars über dem Kapiteltitel wurden als "zu viel Information"
    empfunden — der Kapiteltitel allein reicht für Section-Orientierung.

    Args:
        slide: Slide-Objekt (existierende Folie!)
        sections: list[str] — Section-Namen
        current_idx: int (0-basiert) — welche Section aktiv ist
        show_breadcrumb: False (Default) — keine Breadcrumb-Bars zeigen.
            True → Legacy-Verhalten mit horizontalen Bar-Segmenten + Labels oben.
        y_pos: y-Koordinate für Breadcrumb (nur wenn show_breadcrumb=True)
        override_kicker: True (Default) — Kapiteltitel der Folie wird mit dem
            aktiven Section-Namen überschrieben.

    Empfohlene Nutzung (v0.1.2+):
        s = chart_bar(prs, ..., kicker="", headline="Umsatz 2025: Print dominant")
        add_section_tracker(s, ["Ausgangslage", "Markt", ...], current_idx=0)
        # → Kapiteltitel zeigt "Ausgangslage", keine Bars oben.
    """
    n = len(sections)
    if n == 0: return slide

    # Optional: Breadcrumb-Bars (Legacy, default OFF)
    if show_breadcrumb:
        if y_pos is None: y_pos = Inches(0.40)
        seg_gap = int(Inches(0.10))
        seg_w = (int(GW) - (n - 1) * seg_gap) // n
        bar_h = int(Inches(0.08))

        for i, section in enumerate(sections):
            x = int(GX) + i * (seg_w + seg_gap)
            is_active = (i == current_idx)
            is_past = (i < current_idx)
            bar_color = RED if is_active else (G2 if is_past else G3)
            bar_thickness = bar_h if is_active else int(bar_h * 0.6)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x, int(y_pos), seg_w, bar_thickness)
            bar.fill.solid(); bar.fill.fore_color.rgb = bar_color
            bar.line.fill.background(); bar.shadow.inherit = False
            label_color = BLACK if is_active else G2
            label_font = SB if is_active else BODY
            tf = tb(slide, x, int(y_pos + Inches(0.14)), seg_w, int(Inches(0.30)))
            run(tf.paragraphs[0], section, 10, label_font, label_color, space_after=0)

    # Kapiteltitel mit aktivem Section-Namen überschreiben (Default-Verhalten)
    if override_kicker and 0 <= current_idx < n:
        active_name = sections[current_idx]
        # Weisses Rechteck deckt vorhandenen Kapiteltitel-Textbox ab
        cover = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        int(GX), int(GY_K),
                                        int(GW), int(Inches(0.42)))
        cover.fill.solid(); cover.fill.fore_color.rgb = WHITE
        cover.line.fill.background(); cover.shadow.inherit = False
        # Neuer Kapiteltitel = Section-Name
        ktf = tb(slide, GX, GY_K, GW, Inches(0.40), MSO_ANCHOR.TOP)
        run(ktf.paragraphs[0], active_name, KICK_PT, BODY, BLACK, space_after=0)

    return slide

def scqa_structure(prs, situation, complication, question, answer,
                    kicker="Storyline", folio_start=None):
    """SCQA Pitch-Struktur (Situation → Complication → Question → Answer).

    Erzeugt 4 Folien-Block, der die klassische MBB-Pitch-Struktur abbildet.
    Jede Folie erhält ein Tracker-Element, das den SCQA-Fortschritt zeigt.

    Args:
        situation, complication, question, answer:
            jeweils dict mit "headline" und "body" — body ist Multi-line String
            mit '· '-Bullets (siehe add_content body_text format)
            ODER tuple(headline, body)
        kicker: Kapiteltitel-Text auf allen 4 Folien (Default "Storyline")
        folio_start: optional, z.B. "3" → Folien 3/N, 4/N, 5/N, 6/N

    Beispiel:
        scqa_structure(prs,
            situation={
                "headline": "Print dominiert noch",
                "body": "· 65% Umsatz aus Print\\n· 5 Verlagsobjekte etabliert\\n· Stabile Bestandskunden"},
            complication={
                "headline": "Print-Markt schrumpft",
                "body": "· -8% Werbeerlöse 2026\\n· -7% Auflagen YoY\\n· Cross-Subventionierung nicht nachhaltig"},
            question={
                "headline": "Wie reagieren?",
                "body": "· Digital-Investitionen verstärken?\\n· Print-Portfolio optimieren?\\n· Neue Geschäftsmodelle?"},
            answer={
                "headline": "3-Säulen-Strategie",
                "body": "· Print profitabel halten\\n· Digital aggressiv ausbauen\\n· AI/Plattform als Horizont 3"})
    """
    phases = [
        ("S", "Situation",    situation,    G2),
        ("C", "Complication", complication, RED),
        ("Q", "Question",     question,     G1),
        ("A", "Answer",       answer,       TURK),
    ]
    section_names = ["Situation", "Complication", "Question", "Answer"]

    slides_created = []
    for i, (code, label, data, accent) in enumerate(phases):
        # Daten extrahieren
        if isinstance(data, dict):
            headline = data.get("headline", "")
            body = data.get("body", "")
        elif isinstance(data, tuple) and len(data) == 2:
            headline = data[0]; body = data[1]
        else:
            headline = ""; body = ""

        # Folio berechnen falls start gesetzt
        if folio_start is not None:
            try:
                folio_num = int(folio_start) + i
                folio = str(folio_num)
            except (ValueError, TypeError):
                folio = ""
        else:
            folio = ""

        # Inhaltsfolie via add_content
        s = add_content(prs, "viel", kicker, headline, body, folio=folio)

        # SCQA-Tracker oben (verwendet add_section_tracker)
        add_section_tracker(s, section_names, current_idx=i, y_pos=Inches(0.30))

        # Code-Badge oben rechts (gross, in Accent-Farbe)
        badge_size = int(Inches(0.85))
        badge_x = int(GX + GW - badge_size - Inches(0.20))
        badge_y = int(Inches(0.85))
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    badge_x, badge_y, badge_size, badge_size)
        badge.fill.solid(); badge.fill.fore_color.rgb = accent
        badge.line.color.rgb = WHITE; badge.line.width = Pt(2)
        badge.shadow.inherit = False
        bf = tb(s, badge_x, badge_y, badge_size, badge_size, MSO_ANCHOR.MIDDLE)
        run(bf.paragraphs[0], code, 40, SB, WHITE, align=PP_ALIGN.CENTER, space_after=0)

        slides_created.append(s)

    return slides_created

# ═════════════════════════════════════════════════════════════════════════════
# v0.2.0 STUBS — Phase 2 Full MBB-Level (Sprint 8–13)
# Optional, nach Pilot-User-Feedback aus Phase 1.
# ═════════════════════════════════════════════════════════════════════════════

# ── Sprint 8: Erweiterte Frameworks (v0.2.0-alpha.1 — IMPLEMENTED) ───────────

def mckinsey_9box(prs, units, kicker="", headline="", folio="", source=""):
    """McKinsey 9-Box GE-Matrix — Marktattraktivität × Wettbewerbsstärke.

    Args:
        units: list[dict] mit:
            "name":           str
            "attractiveness": float [0,1] — Marktattraktivität (1=hoch=oben)
            "strength":       float [0,1] — Wettbewerbsstärke (1=hoch=LINKS, klassisch invertiert)
            "revenue":        float (optional) — Umsatz, bestimmt Bubble-Grösse

    9 Felder mit Strategie-Empfehlungen:
        Top-L: Investieren | Top-M: Selektiv investieren | Top-R: Selektiv
        Mid-L: Selektiv invest. | Mid-M: Selektiv | Mid-R: Ernten/Desinvest.
        Bot-L: Selektiv | Bot-M: Ernten/Desinvest. | Bot-R: Desinvestieren

    Beispiel:
        mckinsey_9box(prs, units=[
            {"name":"Fachmedien Print", "attractiveness":0.40, "strength":0.85, "revenue":1200},
            {"name":"Digital-Portal",   "attractiveness":0.85, "strength":0.35, "revenue":450},
            {"name":"Events",           "attractiveness":0.65, "strength":0.55, "revenue":380},
        ], kicker="Portfolio", headline="9-Box Galledia 2026")
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    # Grid: 9 Felder, 3x3
    qx = _CHART_LEFT + Inches(1.4)  # links Platz für Y-Achse-Label
    qy = _CHART_TOP
    qw = _CHART_WIDTH - Inches(2.0)
    qh = _CHART_HEIGHT - Inches(0.5)
    third_w = qw // 3
    third_h = qh // 3

    # Strategie-Empfehlungen + Farbtönung (Investieren=TURK, Selektiv=GL, Desinvest=G3)
    GRID_RECS = [
        # Row 0 (top): high attractiveness
        ("Investieren", TURK), ("Selektiv investieren", TURK), ("Selektiv", GL),
        # Row 1 (mid)
        ("Selektiv investieren", TURK), ("Selektiv", GL), ("Ernten / Desinvest.", G3),
        # Row 2 (bot): low attractiveness
        ("Selektiv", GL), ("Ernten / Desinvest.", G3), ("Desinvestieren", G3),
    ]

    for row in range(3):
        for col in range(3):
            idx = row * 3 + col
            rec, color = GRID_RECS[idx]
            # Sehr leichter Tönung — TURK light, G3 light
            if color == TURK:
                fill = RGBColor(0xCC, 0xEE, 0xEC)
            elif color == G3:
                fill = RGBColor(0xE6, 0xE6, 0xE6)
            else:
                fill = GL
            x_ = qx + col * third_w
            y_ = qy + row * third_h
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x_), int(y_), int(third_w), int(third_h))
            rect.fill.solid(); rect.fill.fore_color.rgb = fill
            rect.line.color.rgb = G4; rect.line.width = Pt(0.75)
            rect.shadow.inherit = False
            # Empfehlung im Feld (oben links, dezent)
            rec_tf = tb(s, int(x_ + Inches(0.15)), int(y_ + Inches(0.10)),
                        int(third_w - Inches(0.3)), Inches(0.35))
            run(rec_tf.paragraphs[0], rec, 10, BODY, G2, italic=True, space_after=0)

    # Achsen-Labels
    # X: hoch (links) bis niedrig (rechts) — INVERTIERT wie BCG
    xl = tb(s, int(qx), int(qy + qh + Inches(0.12)), int(qw // 2), Inches(0.30))
    run(xl.paragraphs[0], "← Hohe Wettbewerbsstärke", 11, BODY, G2, align=PP_ALIGN.LEFT, space_after=0)
    xr = tb(s, int(qx + qw // 2), int(qy + qh + Inches(0.12)), int(qw // 2), Inches(0.30))
    run(xr.paragraphs[0], "Niedrige Wettbewerbsstärke →", 11, BODY, G2, align=PP_ALIGN.RIGHT, space_after=0)
    # Y: hoch (oben) bis niedrig (unten)
    yt = tb(s, int(qx - Inches(1.35)), int(qy + Inches(0.10)), Inches(1.25), Inches(0.30))
    run(yt.paragraphs[0], "Hohe Attraktivität ↑", 11, BODY, G2, align=PP_ALIGN.RIGHT, space_after=0)
    yb = tb(s, int(qx - Inches(1.35)), int(qy + qh - Inches(0.40)), Inches(1.25), Inches(0.30))
    run(yb.paragraphs[0], "Niedrige Attraktivität ↓", 11, BODY, G2, align=PP_ALIGN.RIGHT, space_after=0)

    # Bubbles (klassisch: hohe strength = LINKS, hohe attractiveness = OBEN)
    max_rev = max((u.get("revenue", 1) for u in units), default=1) or 1
    for u in units:
        attr = max(0.0, min(1.0, u.get("attractiveness", 0.5)))
        stre = max(0.0, min(1.0, u.get("strength", 0.5)))
        rev  = u.get("revenue", 1)
        name = u.get("name", "?")
        bubble_d = int(Inches(0.40 + 0.80 * (rev / max_rev)))

        # X: strength=1 → ix=0 (links). Y: attractiveness=1 → iy=0 (oben).
        ix = 1 - stre
        iy = 1 - attr
        px = int(qx + ix * qw) - bubble_d // 2
        py = int(qy + iy * qh) - bubble_d // 2

        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, px, py, bubble_d, bubble_d)
        circle.fill.solid(); circle.fill.fore_color.rgb = RED
        circle.line.color.rgb = WHITE; circle.line.width = Pt(2)
        circle.shadow.inherit = False
        # Label rechts neben Bubble
        ltf = tb(s, px + bubble_d + int(Inches(0.05)), py - int(Inches(0.02)),
                 int(Inches(3.0)), int(Inches(0.4)))
        run(ltf.paragraphs[0], name, 11, SB, BLACK, align=PP_ALIGN.LEFT, space_after=0)

    _set_footer(s, prs, folio, source); return s

def bowman_strategy_clock(prs, position, competitors=None, kicker="", headline="", folio="", source=""):
    """Bowman Strategy Clock — Preis × Mehrwert, 8 strategische Positionen auf einem Kreis.

    Args:
        position: int 1-8 — eigene Position (Galledia, in RED markiert):
            1 = Low Price, No Frills
            2 = Low Price
            3 = Hybrid (preis-leistung)
            4 = Differentiation
            5 = Focused Differentiation
            6 = Risky High Margin
            7 = Monopoly Pricing
            8 = Loss of Share
        competitors: optional list[tuple[name, position_int]] — Wettbewerber-Positionen

    Layout: Kreis (Clock) mit 8 Labels am Umfang, RED-Marker für eigene Position,
            graue Marker für Wettbewerber.
    """
    import math
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    cx = _CHART_LEFT + _CHART_WIDTH // 2
    cy = _CHART_TOP + _CHART_HEIGHT // 2
    radius = int(min(_CHART_WIDTH, _CHART_HEIGHT) * 0.36)

    # Kreis (transparenter Rand)
    inner_d = int(radius * 0.5)
    inner_circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                        int(cx - radius), int(cy - radius),
                                        int(radius * 2), int(radius * 2))
    inner_circle.fill.solid(); inner_circle.fill.fore_color.rgb = GL
    inner_circle.line.color.rgb = G3; inner_circle.line.width = Pt(1.5)
    inner_circle.shadow.inherit = False

    # 8 Positionen, beginnend "Low Price No Frills" links unten (≈7.5 Uhr / 225°)
    # Rotation gegen den Uhrzeigersinn (CCW): 1=225°, 2=270°(unten), 3=315°(rechts unten),
    #   4=0°(rechts), 5=45°(rechts oben), 6=90°(oben), 7=135°(links oben), 8=180°(links)
    POSITIONS = [
        ("No Frills",              "Niedriger Preis, geringer Mehrwert"),
        ("Low Price",              "Niedrigpreis-Strategie"),
        ("Hybrid",                 "Preis-Leistungs-Sieger"),
        ("Differentiation",        "Hoher Mehrwert, mittlerer Preis"),
        ("Focused Differentiation","Premium-Nische"),
        ("Risky High Margin",      "Hoher Preis, mittlerer Mehrwert"),
        ("Monopoly Pricing",       "Hoher Preis, ohne Mehrwert"),
        ("Loss of Share",          "Mittlerer Preis, niedriger Mehrwert"),
    ]
    angles_deg = [225, 270, 315, 0, 45, 90, 135, 180]

    # Marker-Größe
    marker_d = int(Inches(0.42))

    for i, ((name, desc), ang_deg) in enumerate(zip(POSITIONS, angles_deg)):
        ang = math.radians(ang_deg)
        # Marker-Position am Umfang
        mx = cx + int(radius * math.cos(ang)) - marker_d // 2
        my = cy + int(radius * math.sin(ang)) - marker_d // 2
        idx = i + 1
        is_self = (idx == position)
        is_competitor = competitors and any(p == idx for _, p in competitors)
        marker_color = RED if is_self else (G2 if is_competitor else G3)
        marker = s.shapes.add_shape(MSO_SHAPE.OVAL, mx, my, marker_d, marker_d)
        marker.fill.solid(); marker.fill.fore_color.rgb = marker_color
        marker.line.color.rgb = WHITE; marker.line.width = Pt(2)
        marker.shadow.inherit = False
        # Position-Nummer im Marker
        ntf = tb(s, mx, my, marker_d, marker_d, MSO_ANCHOR.MIDDLE)
        run(ntf.paragraphs[0], str(idx), 14, SB, WHITE, align=PP_ALIGN.CENTER, space_after=0)

        # Label aussen am Marker
        label_offset = int(Inches(0.85))
        lx = cx + int((radius + label_offset) * math.cos(ang))
        ly = cy + int((radius + label_offset) * math.sin(ang))
        # Label-Box etwas größer, zentriert um (lx, ly)
        lw, lh = int(Inches(2.6)), int(Inches(0.55))
        ltf = tb(s, lx - lw // 2, ly - lh // 2, lw, lh, MSO_ANCHOR.MIDDLE)
        run(ltf.paragraphs[0], f"{idx}. {name}", 11,
            SB if is_self else BODY,
            BLACK if is_self else G1,
            align=PP_ALIGN.CENTER, space_after=2)
        run(ltf.add_paragraph(), desc, 9, BODY, G2, align=PP_ALIGN.CENTER, space_after=0)

    # Achsen-Beschriftung in der Mitte
    axes_tf = tb(s, int(cx - Inches(1.5)), int(cy - Inches(0.3)),
                 int(Inches(3.0)), int(Inches(0.6)), MSO_ANCHOR.MIDDLE)
    run(axes_tf.paragraphs[0], "Preis →", 10, BODY, G2, align=PP_ALIGN.CENTER, space_after=2)
    run(axes_tf.add_paragraph(), "Mehrwert ↑", 10, BODY, G2, align=PP_ALIGN.CENTER, space_after=0)

    _set_footer(s, prs, folio, source); return s

def wardley_map(prs, components, edges=None, kicker="", headline="", folio="", source=""):
    """Wardley Map — Wertschöpfung (y) × Reife (x: Genesis → Commodity).

    Args:
        components: list[dict] mit:
            "name":      str
            "value":     float [0,1] — Wertschöpfungs-Höhe (1=user-facing, oben)
            "evolution": float [0,1] — Reife (0=Genesis, 1=Commodity, links nach rechts)
            "is_us":     optional bool — eigener Komponent (RED Marker, sonst G2)
        edges: optional list[tuple[name_a, name_b]] — Verbindungen zwischen Komponenten

    Layout: X-Achse Reife mit 4 Phasen (Genesis | Custom | Product | Commodity),
            Y-Achse Wertschöpfung, Komponenten als Punkte mit Verbindungslinien.
    """
    import math
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    qx = _CHART_LEFT + Inches(1.2)
    qy = _CHART_TOP + Inches(0.3)
    qw = _CHART_WIDTH - Inches(1.8)
    qh = _CHART_HEIGHT - Inches(1.0)

    # Hintergrund
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(qx), int(qy), int(qw), int(qh))
    bg.fill.solid(); bg.fill.fore_color.rgb = GL
    bg.line.color.rgb = G4; bg.line.width = Pt(0.5)
    bg.shadow.inherit = False

    # 4 vertikale Phasen-Linien
    phases = ["Genesis", "Custom-Built", "Product / Rental", "Commodity / Utility"]
    phase_w = qw // 4
    from pptx.enum.shapes import MSO_CONNECTOR
    for i in range(1, 4):
        sep_x = qx + i * phase_w
        line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(sep_x), int(qy), int(sep_x), int(qy + qh))
        line.line.color.rgb = G4; line.line.width = Pt(0.5)
    # Phase-Labels unter Box
    for i, phase in enumerate(phases):
        lx = qx + i * phase_w
        plt_tf = tb(s, int(lx), int(qy + qh + Inches(0.10)), int(phase_w), Inches(0.30))
        run(plt_tf.paragraphs[0], phase, 11, SB, G1, align=PP_ALIGN.CENTER, space_after=0)
    # X-Achse Pfeil + Label
    xa_tf = tb(s, int(qx), int(qy + qh + Inches(0.45)), int(qw), Inches(0.30))
    run(xa_tf.paragraphs[0], "Reife / Evolution →", 11, BODY, G2, align=PP_ALIGN.CENTER, space_after=0)
    # Y-Achse Label
    ya_tf = tb(s, int(qx - Inches(1.15)), int(qy + qh // 2 - Inches(0.20)),
               Inches(1.05), Inches(0.40))
    run(ya_tf.paragraphs[0], "↑ Wertschöpfung", 11, SB, G1, align=PP_ALIGN.RIGHT, space_after=0)

    # Lookup table für Komponent-Positionen (für edges)
    component_pos = {}
    for comp in components:
        name = comp.get("name", "?")
        val  = max(0.0, min(1.0, comp.get("value", 0.5)))
        evo  = max(0.0, min(1.0, comp.get("evolution", 0.5)))
        px = qx + int(evo * qw)
        py = qy + int((1 - val) * qh)
        component_pos[name] = (px, py)

    # Edges zuerst (damit unter den Knoten)
    if edges:
        for a, b in edges:
            if a in component_pos and b in component_pos:
                xa, ya = component_pos[a]
                xb, yb = component_pos[b]
                line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(xa), int(ya), int(xb), int(yb))
                line.line.color.rgb = G3; line.line.width = Pt(1.2)

    # Komponenten als Knoten
    node_d = int(Inches(0.30))
    for comp in components:
        name = comp.get("name", "?")
        is_us = comp.get("is_us", False)
        cx, cy = component_pos[name]
        node = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    cx - node_d // 2, cy - node_d // 2, node_d, node_d)
        node.fill.solid(); node.fill.fore_color.rgb = RED if is_us else G1
        node.line.color.rgb = WHITE; node.line.width = Pt(2)
        node.shadow.inherit = False
        # Label rechts daneben
        ltf = tb(s, cx + node_d, cy - int(Inches(0.10)),
                 int(Inches(2.5)), int(Inches(0.35)))
        run(ltf.paragraphs[0], name, 11,
            SB if is_us else BODY,
            BLACK if is_us else G1,
            align=PP_ALIGN.LEFT, space_after=0)

    _set_footer(s, prs, folio, source); return s

def business_model_canvas(prs, segments, kicker="", headline="", folio="", source=""):
    """Business Model Canvas nach Osterwalder — 9-Felder-Layout.

    Args:
        segments: dict mit 9 Keys (jeweils list[str] mit Bullet-Items):
            "key_partners"        — Schlüsselpartner
            "key_activities"      — Schlüsselaktivitäten
            "key_resources"       — Schlüsselressourcen
            "value_propositions"  — Wertangebote (Zentrum, in RED)
            "customer_relationships" — Kundenbeziehungen
            "channels"            — Vertriebskanäle
            "customer_segments"   — Kundensegmente
            "cost_structure"      — Kostenstruktur
            "revenue_streams"     — Einnahmequellen

    Layout (5 Spalten oben, 2 unten):
        | Partners | Activities  | Value | Relations | Segments |
        |          | Resources   | Prop  | Channels  |          |
        | Cost Structure                | Revenue Streams        |
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    gx0 = _CHART_LEFT
    gy0 = _CHART_TOP
    gw  = _CHART_WIDTH
    gh  = _CHART_HEIGHT
    gap = Inches(0.10)

    # Top-Section ist ~70%, Bottom-Section (Cost/Revenue) ~30% der Höhe
    top_h = int(gh * 0.70) - int(gap // 2)
    bot_h = int(gh * 0.30) - int(gap // 2)

    # 5 Spalten in der Top-Section. Spalte 0,1,3,4 = full height; Spalte 2 = full height (Value Prop)
    # Spalte 1 (Activities) und Spalte 3 (Relations) sind vertikal halbiert (Resources unten, Channels unten)
    col_w = (gw - 4 * gap) // 5  # 5 Spalten
    half_h = (top_h - gap) // 2

    def _cell(x, y, w, h, title, items, bg=GL, text_color=BLACK, title_color=None):
        title_color = title_color or text_color
        card(s, int(x), int(y), int(w), int(h), bg, corner=0.02)
        ttf = tb(s, int(x + Inches(0.15)), int(y + Inches(0.10)),
                 int(w - Inches(0.3)), Inches(0.32))
        run(ttf.paragraphs[0], title, 11, SB, title_color, space_after=2)
        if items:
            body_tf = tb(s, int(x + Inches(0.15)), int(y + Inches(0.45)),
                         int(w - Inches(0.3)), int(h - Inches(0.55)))
            first = True
            for item in (items or [])[:5]:
                p = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
                run(p, f"•  {item}", 9, BODY, text_color, space_after=2)
                first = False

    # Spalte 0: Key Partners (volle Top-Höhe)
    _cell(gx0, gy0, col_w, top_h, "Key Partners",
          segments.get("key_partners", []))
    # Spalte 1 oben: Key Activities, unten: Key Resources
    _cell(gx0 + col_w + gap, gy0, col_w, half_h, "Key Activities",
          segments.get("key_activities", []))
    _cell(gx0 + col_w + gap, gy0 + half_h + gap, col_w, half_h, "Key Resources",
          segments.get("key_resources", []))
    # Spalte 2: Value Propositions (RED, volle Top-Höhe)
    _cell(gx0 + 2 * (col_w + gap), gy0, col_w, top_h, "Value Propositions",
          segments.get("value_propositions", []), bg=RED, text_color=WHITE, title_color=WHITE)
    # Spalte 3 oben: Customer Relationships, unten: Channels
    _cell(gx0 + 3 * (col_w + gap), gy0, col_w, half_h, "Customer Relationships",
          segments.get("customer_relationships", []))
    _cell(gx0 + 3 * (col_w + gap), gy0 + half_h + gap, col_w, half_h, "Channels",
          segments.get("channels", []))
    # Spalte 4: Customer Segments
    _cell(gx0 + 4 * (col_w + gap), gy0, col_w, top_h, "Customer Segments",
          segments.get("customer_segments", []))

    # Bottom-Section: Cost Structure (50%) + Revenue Streams (50%)
    bot_y = gy0 + top_h + gap
    half_bot_w = (gw - gap) // 2
    _cell(gx0, bot_y, half_bot_w, bot_h, "Cost Structure",
          segments.get("cost_structure", []), bg=RGBColor(0xE6, 0xE6, 0xE6))
    _cell(gx0 + half_bot_w + gap, bot_y, half_bot_w, bot_h, "Revenue Streams",
          segments.get("revenue_streams", []), bg=RGBColor(0xCC, 0xEE, 0xEC))

    _set_footer(s, prs, folio, source); return s

def lewin_change(prs, unfreeze, change, refreeze, kicker="", headline="", folio="", source=""):
    """Lewin 3-Phasen Change-Model — Unfreeze → Change → Refreeze.

    Args:
        unfreeze, change, refreeze: dict{'title','items'} oder tuple(title, items)
            Items als list[str] — Massnahmen pro Phase
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    n = 3
    arrow_w = Inches(0.5)
    box_w = (_CHART_WIDTH - 2 * arrow_w) // n
    box_h = _CHART_HEIGHT - Inches(1.0)
    box_y = _CHART_TOP + Inches(0.5)

    phases = [
        ("1. Unfreeze",  "Auftauen — Veränderungsdruck schaffen", unfreeze,  G2),
        ("2. Change",    "Verändern — neue Strukturen einführen", change,    RED),
        ("3. Refreeze",  "Einfrieren — neue Praxis verankern",     refreeze,  TURK),
    ]

    for i, (num_label, desc, data, accent) in enumerate(phases):
        x = _CHART_LEFT + i * (box_w + arrow_w)
        title, items = _normalize_force(data)

        # Box mit Accent-Farbe oben
        card(s, int(x), int(box_y), int(box_w), int(box_h), GL, corner=0.04)
        # Accent-Banner oben (Phase-Indicator)
        banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      int(x), int(box_y), int(box_w), int(Inches(0.7)))
        banner.fill.solid(); banner.fill.fore_color.rgb = accent
        banner.line.fill.background(); banner.shadow.inherit = False
        # Phase-Nummer + Name im Banner
        btf = tb(s, int(x + Inches(0.20)), int(box_y), int(box_w - Inches(0.4)), int(Inches(0.7)), MSO_ANCHOR.MIDDLE)
        run(btf.paragraphs[0], num_label, 18, SB, WHITE, space_after=0)
        # Beschreibung unter dem Banner
        dtf = tb(s, int(x + Inches(0.30)), int(box_y + Inches(0.80)),
                 int(box_w - Inches(0.6)), Inches(0.40))
        run(dtf.paragraphs[0], desc, 11, BODY, G1, italic=True, space_after=4)
        # User-Title (falls anders als default)
        if title:
            ttf = tb(s, int(x + Inches(0.30)), int(box_y + Inches(1.30)),
                     int(box_w - Inches(0.6)), Inches(0.40))
            run(ttf.paragraphs[0], title, 14, SB, BLACK, space_after=2)
        # Bullet-Items
        if items:
            items_y = box_y + Inches(1.30 if not title else 1.80)
            items_tf = tb(s, int(x + Inches(0.30)), int(items_y),
                          int(box_w - Inches(0.6)), int(box_h - (items_y - box_y) - Inches(0.30)))
            first = True
            for item in items[:6]:
                p = items_tf.paragraphs[0] if first else items_tf.add_paragraph()
                run(p, f"•  {item}", 12, BODY, BLACK, space_after=4)
                first = False

        # Pfeil zur nächsten Phase
        if i < n - 1:
            arrow_x = x + box_w
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        int(arrow_x), int(box_y + box_h // 2 - Inches(0.30)),
                                        int(arrow_w), int(Inches(0.6)))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = G3
            arrow.line.fill.background(); arrow.shadow.inherit = False

    _set_footer(s, prs, folio, source); return s

def kotter_8step(prs, steps, current_step=None, kicker="", headline="", folio="", source=""):
    """Kotter's 8-Step Change-Process.

    Args:
        steps: list[dict] mit "title" und optional "items" (max 8 Schritte)
            ODER list[tuple[title, items]]
            ODER list[str] (nur Titles)
            Standard-Reihenfolge (wenn weniger als 8 gegeben):
              1. Dringlichkeit erzeugen
              2. Führungskoalition bilden
              3. Vision entwickeln
              4. Vision kommunizieren
              5. Mitarbeitende befähigen
              6. Kurzfristige Erfolge sichern
              7. Veränderung konsolidieren
              8. Neue Kultur verankern
        current_step: int 1-8 — aktueller Schritt (RED-Highlight)

    Layout: 2 Reihen à 4 Schritte (oder 1 Reihe wenn ≤4), nummerierte Boxen.
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    DEFAULT_TITLES = [
        "Dringlichkeit erzeugen",
        "Führungskoalition bilden",
        "Vision entwickeln",
        "Vision kommunizieren",
        "Mitarbeitende befähigen",
        "Kurzfristige Erfolge sichern",
        "Veränderung konsolidieren",
        "Neue Kultur verankern",
    ]

    # Steps normalisieren
    normalized = []
    for i, step in enumerate(steps[:8]):
        if isinstance(step, dict):
            title = step.get("title", DEFAULT_TITLES[i] if i < 8 else "")
            items = step.get("items", [])
        elif isinstance(step, (tuple, list)) and len(step) == 2:
            title = step[0]; items = list(step[1]) if isinstance(step[1], (list, tuple)) else []
        else:
            title = str(step); items = []
        normalized.append((title, items))
    # Fill up to 8 mit Defaults
    while len(normalized) < 8:
        normalized.append((DEFAULT_TITLES[len(normalized)], []))

    n = 8
    cols = 4
    rows = 2
    gap = Inches(0.15)
    box_w = (_CHART_WIDTH - (cols - 1) * gap) // cols
    box_h = (_CHART_HEIGHT - (rows - 1) * gap) // rows

    for i, (title, items) in enumerate(normalized):
        row = i // cols
        col = i % cols
        x = _CHART_LEFT + col * (box_w + gap)
        y = _CHART_TOP + row * (box_h + gap)
        is_current = (current_step is not None and (i + 1) == current_step)

        # Box
        bg = RED if is_current else GL
        text_color = WHITE if is_current else BLACK
        card(s, int(x), int(y), int(box_w), int(box_h), bg, corner=0.04)

        # Step-Nummer als Badge oben links
        badge_d = int(Inches(0.65))
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    int(x + Inches(0.15)), int(y + Inches(0.15)),
                                    badge_d, badge_d)
        badge.fill.solid(); badge.fill.fore_color.rgb = WHITE if is_current else RED
        badge.line.fill.background(); badge.shadow.inherit = False
        btf = tb(s, int(x + Inches(0.15)), int(y + Inches(0.15)),
                 badge_d, badge_d, MSO_ANCHOR.MIDDLE)
        run(btf.paragraphs[0], str(i + 1), 22, SB,
            RED if is_current else WHITE, align=PP_ALIGN.CENTER, space_after=0)

        # Title rechts neben Badge
        ttf = tb(s, int(x + Inches(0.90)), int(y + Inches(0.20)),
                 int(box_w - Inches(1.05)), Inches(0.65))
        run(ttf.paragraphs[0], title, 12, SB, text_color, space_after=2)

        # Items unter Title
        if items:
            items_tf = tb(s, int(x + Inches(0.20)), int(y + Inches(0.95)),
                          int(box_w - Inches(0.4)), int(box_h - Inches(1.10)))
            first = True
            for item in items[:3]:
                p = items_tf.paragraphs[0] if first else items_tf.add_paragraph()
                run(p, f"•  {item}", 10, BODY, text_color, space_after=2)
                first = False

    _set_footer(s, prs, folio, source); return s

# ── Sprint 9: Advanced Charts (v0.2.0-alpha.2 — IMPLEMENTED) ─────────────────

def lever_tornado(prs, levers, kicker="", headline="", folio="", source="",
                   x_label="Impact"):
    """Tornado-Diagramm für Sensitivitäts-Analyse — Bars sortiert nach Impact-Range.

    Args:
        levers: list[tuple[label, low, high]] — Hebel mit unterer und oberer Schwelle
            Beispiel: [("Preis +/-10%", -240, 280), ("Kosten +/-5%", -120, 180), ...]
        x_label: Achsen-Beschriftung (Default "Impact")

    Layout: Vertikale Linie bei 0, horizontale Bars zu beiden Seiten,
            sortiert nach |range| absteigend (höchster Impact oben).
            Negative Werte in RED, positive in TURK.
    """
    plt = _init_matplotlib()
    # Sortierung nach absoluter Range (high - low) absteigend
    levers_sorted = sorted(levers, key=lambda l: abs(l[2] - l[1]), reverse=True)
    labels = [l[0] for l in levers_sorted]
    lows  = [l[1] for l in levers_sorted]
    highs = [l[2] for l in levers_sorted]

    fig, ax = plt.subplots(figsize=(18, 6.26))
    y_pos = list(range(len(labels)))
    # Negative Bars (lows)
    ax.barh(y_pos, lows, color=_GALLEDIA_HEX["RED"], edgecolor="white", linewidth=1, height=0.7)
    # Positive Bars (highs)
    ax.barh(y_pos, highs, color=_GALLEDIA_HEX["TURK"], edgecolor="white", linewidth=1, height=0.7)
    # Zero-Line
    ax.axvline(0, color=_GALLEDIA_HEX["G1"], linewidth=1.5)
    # Labels
    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=11)
    ax.invert_yaxis()  # Höchster Impact oben
    ax.set_xlabel(x_label)
    # Wert-Labels an Bar-Enden
    for i, (lo, hi) in enumerate(zip(lows, highs)):
        if lo != 0:
            ax.text(lo, i, f" {lo:+,}".replace(",", "'"), va="center", ha="right" if lo < 0 else "left",
                    fontsize=10, color=_GALLEDIA_HEX["RED"], fontweight="bold")
        if hi != 0:
            ax.text(hi, i, f" {hi:+,}".replace(",", "'"), va="center", ha="left" if hi > 0 else "right",
                    fontsize=10, color=_GALLEDIA_HEX["TURK"], fontweight="bold")

    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_tornado_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

def sankey(prs, flows, kicker="", headline="", folio="", source=""):
    """Vereinfachtes Sankey-Diagramm via Shapes — Flows von Quellen zu Zielen.

    Args:
        flows: list[tuple[source, target, value]] — z.B. [("Umsatz", "Personal", 2400), ...]
            Quellen werden links, Ziele rechts platziert. Verbindungen via dünne Linien
            mit Dicke proportional zum Wert.

    Layout: zwei vertikale Spalten (Quellen | Ziele), Verbindungslinien zwischen.
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    # Aggregiere Quellen und Ziele
    sources = {}
    targets = {}
    for src, tgt, val in flows:
        sources[src] = sources.get(src, 0) + val
        targets[tgt] = targets.get(tgt, 0) + val
    src_total = sum(sources.values()) or 1
    tgt_total = sum(targets.values()) or 1

    # Layout: Quellen links, Ziele rechts
    bar_w = Inches(2.5)
    chart_pad_y = Inches(0.3)
    left_x = _CHART_LEFT + Inches(0.5)
    right_x = _CHART_LEFT + _CHART_WIDTH - Inches(0.5) - bar_w
    chart_top = _CHART_TOP + chart_pad_y
    chart_h = _CHART_HEIGHT - 2 * chart_pad_y

    # Y-Positionen für Quellen
    src_positions = {}  # name → (y_top, height)
    y_cursor = chart_top
    src_gap = int(Inches(0.05))
    for name, val in sources.items():
        h = int(chart_h * val / src_total) - src_gap
        src_positions[name] = (y_cursor, h)
        # Source-Bar
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    int(left_x), int(y_cursor), int(bar_w), h)
        rect.fill.solid(); rect.fill.fore_color.rgb = RED
        rect.line.color.rgb = WHITE; rect.line.width = Pt(1)
        rect.shadow.inherit = False
        # Label im/auf der Bar
        if h > int(Inches(0.4)):
            ltf = tb(s, int(left_x + Inches(0.15)), y_cursor, int(bar_w - Inches(0.3)), h, MSO_ANCHOR.MIDDLE)
            run(ltf.paragraphs[0], name, 13, SB, WHITE, space_after=2)
            run(ltf.add_paragraph(), _fmt_num(val), 11, BODY, WHITE, space_after=0)
        y_cursor += h + src_gap

    # Y-Positionen für Ziele
    tgt_positions = {}
    y_cursor = chart_top
    tgt_gap = int(Inches(0.05))
    PALETTE = [G1, G2, TURK, G3, G4]
    for i, (name, val) in enumerate(targets.items()):
        h = int(chart_h * val / tgt_total) - tgt_gap
        tgt_positions[name] = (y_cursor, h)
        color = PALETTE[i % len(PALETTE)]
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    int(right_x), int(y_cursor), int(bar_w), h)
        rect.fill.solid(); rect.fill.fore_color.rgb = color
        rect.line.color.rgb = WHITE; rect.line.width = Pt(1)
        rect.shadow.inherit = False
        if h > int(Inches(0.4)):
            ltf = tb(s, int(right_x + Inches(0.15)), y_cursor, int(bar_w - Inches(0.3)), h, MSO_ANCHOR.MIDDLE)
            text_color = WHITE if color in (G1, G2, TURK) else BLACK
            run(ltf.paragraphs[0], name, 13, SB, text_color, space_after=2)
            run(ltf.add_paragraph(), _fmt_num(val), 11, BODY, text_color, space_after=0)
        y_cursor += h + tgt_gap

    # Flow-Linien zwischen Quellen und Zielen
    from pptx.enum.shapes import MSO_CONNECTOR
    # Track current y-offset in each source/target for stacking flows
    src_offset = {n: 0 for n in sources}
    tgt_offset = {n: 0 for n in targets}
    for src, tgt, val in flows:
        if src not in src_positions or tgt not in tgt_positions: continue
        src_y, src_h = src_positions[src]
        tgt_y, tgt_h = tgt_positions[tgt]
        flow_h_src = int(src_h * val / sources[src])
        flow_h_tgt = int(tgt_h * val / targets[tgt])
        # Start- und End-Mittelpunkte
        y1 = src_y + src_offset[src] + flow_h_src // 2
        y2 = tgt_y + tgt_offset[tgt] + flow_h_tgt // 2
        # Verbindungslinie (gerade, Dicke proportional zu val)
        x1 = int(left_x + bar_w)
        x2 = int(right_x)
        # Use line shape with thickness
        line_thickness = max(0.5, min(8, (val / max(src_total, tgt_total)) * 25))
        line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = G3
        line.line.width = Pt(line_thickness)
        src_offset[src] += flow_h_src
        tgt_offset[tgt] += flow_h_tgt

    # Spalten-Titel
    src_lbl = tb(s, int(left_x), int(chart_top - Inches(0.40)), int(bar_w), Inches(0.30))
    run(src_lbl.paragraphs[0], "Quellen", 12, SB, G2, align=PP_ALIGN.CENTER, space_after=0)
    tgt_lbl = tb(s, int(right_x), int(chart_top - Inches(0.40)), int(bar_w), Inches(0.30))
    run(tgt_lbl.paragraphs[0], "Ziele", 12, SB, G2, align=PP_ALIGN.CENTER, space_after=0)

    _set_footer(s, prs, folio, source); return s

def indexed_comparison(prs, series, base_year, kicker="", headline="", folio="", source="",
                        y_label="Index (Basis = 100)"):
    """Indexierte Trendlinien — alle Serien starten bei 100 im base_year.

    Args:
        series: dict[label, list[values]] — absolute Werte
        base_year: str oder int — Bezeichnung des Basis-Jahrs (für Achsen-Beschriftung)
        y_label: Y-Achsen-Label

    Beispiel:
        indexed_comparison(prs,
            series={"Galledia": [2400, 2520, 2688, 2832], "Markt": [2400, 2448, 2496, 2568]},
            base_year=2023,
            kicker="Wachstum 2023-2026e", headline="Galledia über Markt")
    """
    plt = _init_matplotlib()
    fig, ax = plt.subplots(figsize=(18, 6.26))

    # Normalisierung: erster Wert pro Serie = 100
    x_labels = None
    for i, (name, values) in enumerate(series.items()):
        if not values: continue
        base_val = values[0]
        indexed = [100 * v / base_val for v in values]
        if x_labels is None:
            # Generiere Year-Labels ausgehend von base_year
            try:
                bs = int(base_year)
                x_labels = [str(bs + j) for j in range(len(values))]
            except (ValueError, TypeError):
                x_labels = [f"{base_year}+{j}" if j > 0 else str(base_year) for j in range(len(values))]
        color = _GALLEDIA_PALETTE[i % len(_GALLEDIA_PALETTE)]
        ax.plot(x_labels, indexed, label=name, color=color, linewidth=2.5,
                marker="o", markersize=7, markeredgecolor="white", markeredgewidth=1.5)
        # End-Wert-Annotation
        ax.text(len(x_labels) - 1, indexed[-1], f" {indexed[-1]:.0f}",
                fontsize=11, color=color, fontweight="bold", va="center")

    # Baseline 100
    ax.axhline(100, color=_GALLEDIA_HEX["G3"], linewidth=1, linestyle="--")
    ax.text(0, 100, " Basis = 100", fontsize=10, color=_GALLEDIA_HEX["G2"], va="bottom", ha="left")
    ax.set_ylabel(y_label)
    ax.legend(loc="upper left", frameon=False, fontsize=11)

    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_indexed_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

def trend_with_confidence(prs, x, mean, lower, upper, kicker="", headline="", folio="", source="",
                          y_label="", x_label="", series_label="Mean"):
    """Trend-Linie mit Konfidenz-Band (für Prognosen/Schätzungen).

    Args:
        x: list — x-Achsen-Werte (numerisch oder strings)
        mean: list — Mittelwerte (Hauptlinie)
        lower: list — untere Konfidenzgrenze
        upper: list — obere Konfidenzgrenze
        series_label: Label für die Linie (z.B. "Forecast")

    Layout: Hauptlinie in RED, schattiertes Band zwischen lower und upper in TURK light.
    """
    plt = _init_matplotlib()
    fig, ax = plt.subplots(figsize=(18, 6.26))

    # Konfidenz-Band
    ax.fill_between(x, lower, upper, alpha=0.20, color=_GALLEDIA_HEX["TURK"],
                    label="Konfidenz-Band")
    # Lower/Upper Linien dünn
    ax.plot(x, lower, color=_GALLEDIA_HEX["G3"], linewidth=1, linestyle="--", alpha=0.7)
    ax.plot(x, upper, color=_GALLEDIA_HEX["G3"], linewidth=1, linestyle="--", alpha=0.7)
    # Mean-Linie prominent
    ax.plot(x, mean, label=series_label, color=_GALLEDIA_HEX["RED"], linewidth=2.5,
            marker="o", markersize=7, markeredgecolor="white", markeredgewidth=1.5)

    if y_label: ax.set_ylabel(y_label)
    if x_label: ax.set_xlabel(x_label)
    ax.legend(loc="upper left", frameon=False, fontsize=11)

    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_trend_conf_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

def network_diagram(prs, nodes, edges, kicker="", headline="", folio="", source=""):
    """Stakeholder-Map / System-Architektur als Network-Diagram via Shapes.

    Args:
        nodes: list[dict] mit:
            "name":  str
            "x":     float [0,1] — relative Position horizontal
            "y":     float [0,1] — relative Position vertikal (0=bottom, 1=top)
            "size":  float (optional) — Knoten-Grösse-Multiplikator (Default 1.0)
            "is_us": bool (optional) — eigener Knoten (RED, sonst G2)
            "group": str (optional) — Gruppenname, beeinflusst Farbe
        edges: list[tuple[name_a, name_b]] ODER list[tuple[name_a, name_b, weight]]

    Layout: Knoten an angegebenen Positionen, Edges als Linien (Dicke ~ weight).
    """
    s = _blank(prs)
    if kicker or headline: _set_header(s, kicker, headline)

    qx = _CHART_LEFT + Inches(0.5)
    qy = _CHART_TOP + Inches(0.3)
    qw = _CHART_WIDTH - Inches(1.0)
    qh = _CHART_HEIGHT - Inches(0.6)

    # Hintergrund
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(qx), int(qy), int(qw), int(qh))
    bg.fill.solid(); bg.fill.fore_color.rgb = GL
    bg.line.color.rgb = G4; bg.line.width = Pt(0.5)
    bg.shadow.inherit = False

    # Compute node positions
    node_pos = {}
    for n in nodes:
        name = n.get("name", "?")
        nx = max(0, min(1, n.get("x", 0.5)))
        ny = max(0, min(1, n.get("y", 0.5)))
        px = qx + int(nx * qw)
        py = qy + int((1 - ny) * qh)
        node_pos[name] = (px, py)

    # Edges (Linien zuerst, damit unter den Knoten)
    from pptx.enum.shapes import MSO_CONNECTOR
    for edge in edges:
        if len(edge) == 3:
            a, b, weight = edge
        else:
            a, b = edge[0], edge[1]; weight = 1.0
        if a not in node_pos or b not in node_pos: continue
        xa, ya = node_pos[a]
        xb, yb = node_pos[b]
        line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(xa), int(ya), int(xb), int(yb))
        line.line.color.rgb = G3
        line.line.width = Pt(max(0.5, min(5, weight * 2)))

    # Knoten + Labels
    GROUP_COLORS = {}  # Auto-Cache pro Gruppe
    GROUP_PALETTE = [G1, G2, TURK, G3, RGBColor(0x88, 0x55, 0xAA)]
    for n in nodes:
        name = n.get("name", "?")
        size_mult = n.get("size", 1.0)
        is_us = n.get("is_us", False)
        group = n.get("group", "")
        # Farbe bestimmen
        if is_us:
            color = RED
        elif group:
            if group not in GROUP_COLORS:
                GROUP_COLORS[group] = GROUP_PALETTE[len(GROUP_COLORS) % len(GROUP_PALETTE)]
            color = GROUP_COLORS[group]
        else:
            color = G2

        node_d = int(Inches(0.45 * size_mult))
        px, py = node_pos[name]
        node = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    px - node_d // 2, py - node_d // 2, node_d, node_d)
        node.fill.solid(); node.fill.fore_color.rgb = color
        node.line.color.rgb = WHITE; node.line.width = Pt(2.5)
        node.shadow.inherit = False
        # Label unter dem Knoten
        ltf = tb(s, px - int(Inches(1.0)), py + node_d // 2 + int(Inches(0.05)),
                 int(Inches(2.0)), int(Inches(0.30)))
        run(ltf.paragraphs[0], name, 11,
            SB if is_us else BODY,
            BLACK if is_us else G1,
            align=PP_ALIGN.CENTER, space_after=0)

    _set_footer(s, prs, folio, source); return s

def heatmap_annotated(prs, matrix, row_labels, col_labels, callouts=None,
                      kicker="", headline="", folio="", source="",
                      cmap_low=None, cmap_high=None, value_fmt="{:.0f}"):
    """Heat-Map mit optionalen Callout-Annotations.

    Args:
        matrix: list[list[float]] — n_rows × n_cols Werte-Matrix
        row_labels: list[str] — Zeilen-Beschriftungen
        col_labels: list[str] — Spalten-Beschriftungen
        callouts: optional list[dict] mit "row" (int), "col" (int), "text" (str) —
            Highlight-Pfeile auf bestimmte Zellen
        cmap_low: untere Farbe (Default G4 hell)
        cmap_high: obere Farbe (Default RED)
        value_fmt: Format-String für Zell-Werte

    Layout: Heat-Map via matplotlib imshow mit custom Galledia-Colormap.
    """
    plt = _init_matplotlib()
    import numpy as np
    from matplotlib.colors import LinearSegmentedColormap

    # Galledia-Colormap: Light-Grey → RED
    galledia_cmap = LinearSegmentedColormap.from_list(
        "galledia_red", [
            (1.0, 1.0, 1.0),     # Weiss
            (0.85, 0.85, 0.85),  # Grau
            (0.95, 0.40, 0.55),  # Light Red
            (0.90, 0.11, 0.32),  # Galledia Red
        ], N=256)

    arr = np.array(matrix, dtype=float)
    fig, ax = plt.subplots(figsize=(18, 6.26))
    im = ax.imshow(arr, cmap=galledia_cmap, aspect="auto")

    # Labels
    ax.set_xticks(range(len(col_labels)))
    ax.set_yticks(range(len(row_labels)))
    ax.set_xticklabels(col_labels, fontsize=11)
    ax.set_yticklabels(row_labels, fontsize=11)
    # X-Labels oben statt unten (klassisch Heatmap)
    ax.xaxis.set_label_position("top")
    ax.xaxis.tick_top()

    # Werte in Zellen
    for i in range(len(row_labels)):
        for j in range(len(col_labels)):
            val = arr[i, j]
            # Text-Farbe je nach Hintergrund (helle Zellen → BLACK, dunkle → WHITE)
            normalized = (val - arr.min()) / (arr.max() - arr.min() + 1e-9)
            text_color = "white" if normalized > 0.5 else "black"
            ax.text(j, i, value_fmt.format(val), ha="center", va="center",
                    color=text_color, fontsize=11, fontweight="bold")

    # Colorbar weg (zu viel Info)
    # ax.set_xlabel(""); ax.set_ylabel("")

    # Callouts: Rechtecke um bestimmte Zellen
    if callouts:
        for callout in callouts:
            r = callout.get("row", 0)
            c = callout.get("col", 0)
            text = callout.get("text", "")
            # Rechteck um Zelle
            rect = plt.Rectangle((c - 0.5, r - 0.5), 1, 1, fill=False,
                                  edgecolor=_GALLEDIA_HEX["TURK"], linewidth=3)
            ax.add_patch(rect)
            # Text-Annotation oberhalb der Zelle
            if text:
                ax.annotate(text, xy=(c, r), xytext=(c, r - 1.3),
                            fontsize=10, color=_GALLEDIA_HEX["TURK"], fontweight="bold",
                            ha="center",
                            arrowprops=dict(arrowstyle="->", color=_GALLEDIA_HEX["TURK"], lw=1.5))

    # Gridlines zwischen Zellen
    ax.set_xticks([j - 0.5 for j in range(1, len(col_labels))], minor=True)
    ax.set_yticks([i - 0.5 for i in range(1, len(row_labels))], minor=True)
    ax.grid(which="minor", color="white", linewidth=2)
    ax.tick_params(which="minor", length=0)

    fig.tight_layout()
    png = _render_chart_png(fig, "galledia_heatmap_")
    plt.close(fig)
    return _embed_chart(prs, png, kicker, headline, folio, source)

# ── Sprint 10: Visual Polish (v0.2.0-alpha.3 — IMPLEMENTED) ──────────────────

def _set_shape_alpha(shape, alpha_pct):
    """Setzt Transparenz auf solidFill-Shape. alpha_pct: 0=opaque, 100=transparent."""
    from lxml import etree
    try:
        solidFill = shape.fill._xPr.find(qn('a:solidFill'))
        if solidFill is None: return
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is None: return
        for existing in srgb.findall(qn('a:alpha')):
            srgb.remove(existing)
        # PowerPoint-Alpha-Skala: 0 = fully transparent, 100000 = fully opaque
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', str(int((100 - alpha_pct) * 1000)))
    except Exception:
        pass  # graceful fallback wenn shape-Struktur abweicht

def _download_image(url):
    """Lädt URL → temp file, return local path. Falls schon Pfad: passthrough."""
    if not isinstance(url, str): return url
    if not url.startswith(("http://", "https://")):
        return url
    import urllib.request, tempfile, os
    ext = os.path.splitext(url.split("?")[0])[1] or ".jpg"
    tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False, prefix="galledia_img_")
    tmp.close()
    urllib.request.urlretrieve(url, tmp.name)
    return tmp.name

def add_photo_full(slide, image_path, treatment="none"):
    """Vollflächiges Photo (Slide-Cover) mit optionalem Galledia-CI-Treatment.

    Args:
        slide: Slide-Objekt
        image_path: Lokaler Pfad (PNG/JPG) ODER URL (wird automatisch geladen)
        treatment: 'none' — Bild pur
                   'galledia_red_overlay' — 25% rote Tönung (Galledia-Signatur)
                   'dark_vignette' — Verdunkelung unten (Lesbarkeit für Text-Overlay)
                   'dim' — 40% Schwarz-Overlay (für hellen Text auf dunklem Bild)

    Returns:
        Picture-Shape (für weitere Positionierung)
    """
    path = _download_image(image_path)
    # Vollbild einfügen
    pic = slide.shapes.add_picture(path, 0, 0, width=SH * 16 // 9, height=SH)  # 16:9 mit Folienhöhe
    # Schätzung Folie-Breite: SH (11.25) × 16/9 = 20.0" — passt zur Vorlage_6
    # Korrigieren auf exakte Folienbreite
    slide_w = slide.part.package.presentation_part.presentation.slide_width
    slide_h = slide.part.package.presentation_part.presentation.slide_height
    pic.left = 0; pic.top = 0
    pic.width = slide_w; pic.height = slide_h

    if treatment == "galledia_red_overlay":
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        overlay.fill.solid(); overlay.fill.fore_color.rgb = RED
        overlay.line.fill.background(); overlay.shadow.inherit = False
        _set_shape_alpha(overlay, 75)  # 75% transparent (25% red visible)
    elif treatment == "dim":
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        overlay.fill.solid(); overlay.fill.fore_color.rgb = BLACK
        overlay.line.fill.background(); overlay.shadow.inherit = False
        _set_shape_alpha(overlay, 60)  # 60% transparent (40% black visible)
    elif treatment == "dark_vignette":
        # Verdunkeltes Rechteck nur in der unteren Hälfte (für Text-Lesbarkeit)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          0, slide_h // 2, slide_w, slide_h // 2)
        overlay.fill.solid(); overlay.fill.fore_color.rgb = BLACK
        overlay.line.fill.background(); overlay.shadow.inherit = False
        _set_shape_alpha(overlay, 50)
    return pic

def add_photo_card(slide, x, y, w, h, image_path, caption="", caption_position="bottom"):
    """Eingebettetes Bild mit optionaler Caption (zugeschnitten auf x, y, w, h).

    Args:
        slide: Slide-Objekt
        x, y, w, h: Position und Grösse in EMU
        image_path: Lokaler Pfad ODER URL
        caption: optionale Bildunterschrift (klein, grau)
        caption_position: 'bottom' | 'overlay' (auf Bild) | 'none'

    Returns:
        Picture-Shape
    """
    path = _download_image(image_path)
    # Reserve space für Caption wenn position='bottom'
    if caption and caption_position == "bottom":
        cap_h = int(Inches(0.30))
        pic_h = h - cap_h
    else:
        cap_h = 0
        pic_h = h

    pic = slide.shapes.add_picture(path, int(x), int(y), width=int(w), height=int(pic_h))

    if caption:
        if caption_position == "bottom":
            ctf = tb(slide, int(x), int(y + pic_h), int(w), int(cap_h))
            run(ctf.paragraphs[0], caption, 10, BODY, G2, italic=True, space_after=0)
        elif caption_position == "overlay":
            # Halbtransparentes Schwarz-Banner unten auf dem Bild
            band_h = int(Inches(0.45))
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            int(x), int(y + pic_h - band_h), int(w), int(band_h))
            band.fill.solid(); band.fill.fore_color.rgb = BLACK
            band.line.fill.background(); band.shadow.inherit = False
            _set_shape_alpha(band, 40)  # 60% black visible
            ctf = tb(slide, int(x + Inches(0.15)), int(y + pic_h - band_h), int(w - Inches(0.3)), int(band_h), MSO_ANCHOR.MIDDLE)
            run(ctf.paragraphs[0], caption, 11, BODY, WHITE, space_after=0)
    return pic

def add_icon(slide, icon_path, x, y, size_in=0.6, color=None):
    """Icon einbetten mit optionaler Farbtinting via PIL.

    Args:
        slide: Slide-Objekt
        icon_path: Lokaler Pfad zur Icon-PNG (transparenter Hintergrund empfohlen)
        x, y: Position in EMU
        size_in: Icon-Grösse in Inches (Default 0.6)
        color: 'rot' | 'weiss' | 'schwarz' | None (Original-Farbe)
            Mit Color: PNG wird re-coloriert via PIL (alle nicht-transparenten Pixel)

    Returns:
        Picture-Shape
    """
    final_path = icon_path
    if color and color.lower() in ("rot", "weiss", "schwarz"):
        try:
            from PIL import Image
            import tempfile
            target = {
                "rot":     (0xE6, 0x1C, 0x52),
                "weiss":   (0xFF, 0xFF, 0xFF),
                "schwarz": (0x00, 0x00, 0x00),
            }[color.lower()]
            img = Image.open(icon_path).convert("RGBA")
            data = img.getdata()
            new_data = []
            for r, g, b, a in data:
                if a > 0:
                    new_data.append((*target, a))
                else:
                    new_data.append((r, g, b, a))
            img.putdata(new_data)
            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False, prefix="galledia_icon_")
            tmp.close()
            img.save(tmp.name, "PNG")
            final_path = tmp.name
        except Exception:
            pass  # Fallback: Original-Farbe

    size_emu = int(Inches(size_in))
    pic = slide.shapes.add_picture(final_path, int(x), int(y), width=size_emu, height=size_emu)
    return pic

def simplify_chart(level="default"):
    """Tufte Data-Ink Discipline — passt matplotlib rcParams für minimale Charts an.

    WICHTIG: muss VOR den chart_*-Calls aufgerufen werden — verändert globale
    matplotlib-Defaults für FOLGENDE Charts. Nach Verwendung am besten zurücksetzen
    mit simplify_chart('default').

    Args:
        level: 'default' — Galledia-Standard (light grid, soft colors)
               'light' — weniger Grid, kleinere Ticks
               'medium' — kein Grid, Axes minimal
               'heavy' — pure Daten, nur Datenpunkte und Labels

    Beispiel:
        simplify_chart("medium")
        chart_line(prs, x, series, ...)   # nutzt vereinfachten Stil
        chart_bar(prs, data, ...)         # ebenfalls
        simplify_chart("default")          # zurück zum Standard

    Returns:
        None — modifiziert globale plt.rcParams
    """
    plt = _init_matplotlib()
    if level == "heavy":
        plt.rcParams.update({
            "axes.grid": False,
            "axes.spines.left": False,
            "axes.spines.bottom": False,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "xtick.bottom": False, "ytick.left": False,
            "xtick.major.size": 0, "ytick.major.size": 0,
        })
    elif level == "medium":
        plt.rcParams.update({
            "axes.grid": False,
            "axes.spines.top": False, "axes.spines.right": False,
            "axes.spines.left": True, "axes.spines.bottom": True,
            "xtick.bottom": True, "ytick.left": True,
            "xtick.major.size": 3, "ytick.major.size": 3,
        })
    elif level == "light":
        plt.rcParams.update({
            "axes.grid": True, "grid.alpha": 0.15,
            "axes.spines.top": False, "axes.spines.right": False,
            "xtick.major.size": 4, "ytick.major.size": 4,
        })
    else:  # default
        plt.rcParams.update({
            "axes.grid": True, "grid.alpha": 0.25,
            "grid.color": "#A6A6A6", "grid.linewidth": 0.6,
            "axes.spines.top": False, "axes.spines.right": False,
        })

# ── Sprint 11: Quality-Gates (v0.2.0-alpha.4 — IMPLEMENTED) ──────────────────

# Generische "Worthülsen" die MBB-Berater als Buzzwords vermeiden
_GENERIC_WORDS = {
    "vorteile", "mehrwert", "synergien", "best practice", "best practices",
    "optimierung", "transformation", "effizienz", "innovation", "skalierung",
    "potenzial", "performance-steigerung", "exzellenz", "agilität",
    "next level", "value add",
}

# Catch-All-Phrasen die gegen MECE verstoßen
_CATCH_ALL_PHRASES = ["andere", "etc.", " etc", "und mehr", "weiteres", "u.a.", "u. a."]

# Deutsche Action-Verben (Indikator für action-style Headlines)
_ACTION_VERBS = {
    "wächst", "steigt", "fällt", "sinkt", "spart", "kostet", "erhöht",
    "reduziert", "verdoppelt", "halbiert", "kippt", "dreht", "übersteigt",
    "übertrifft", "unterschreitet", "gewinnt", "verliert", "schafft",
    "wird", "bleibt", "übernimmt", "investiert", "wechselt",
    "schließt", "öffnet", "treibt", "schiebt", "erreicht", "kommt",
    "geht", "läuft", "wirkt", "bringt", "liefert", "stellt", "wagt",
    "startet", "endet", "siegt", "verliert", "expandiert", "schrumpft",
}

def audit_deck(prs, level="warn"):
    """Audit eines kompletten Decks auf MBB-Qualitätskriterien.

    Args:
        prs: Presentation-Objekt
        level: 'warn' — Findings printen + return (default)
               'strict' — bei Findings raise ValueError mit aller Findings-Liste
               'silent' — nur return, keine print-Ausgabe

    Returns:
        list[dict] — jeder Finding mit Keys:
            "slide":      int — Folien-Nummer (1-basiert)
            "severity":   "error" | "warn" | "info"
            "issue":      str — Beschreibung
            "suggestion": str — Verbesserungs-Vorschlag

    Geprüfte Regeln:
        - Generika-Wörter ohne Zahl-Kontext ("Vorteile", "Mehrwert", ...)
        - Bullet-Count: Inhaltsfolien sollten 3-5 Bullets haben (nicht 1-2 = Luftnummer)
        - Headline-Länge (sollte schon durch _check_len validiert sein, redundanter Check)
        - Fehlende Quellen-Angabe (Footnote-Block oder "Quelle:" Text)
        - Folien-Mix: % add_content vs visuelle Layouts

    Beispiel:
        findings = audit_deck(prs, level="warn")
        # → "audit_deck: 3 findings"
        #   "Folie 4 [warn]: Generika-Verdacht: 'Mehrwert' im Folientext"
    """
    findings = []

    n_slides = len(prs.slides)
    add_content_count = 0  # für Mix-Check

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_text = ""
        bullets_count = 0
        has_source_marker = False
        has_kpi_numbers = False
        layout_name = slide.slide_layout.name if hasattr(slide, "slide_layout") else ""

        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text = shape.text_frame.text or ""
            slide_text += text + "\n"
            # Footnote-Indikator
            for sup in "⁰¹²³⁴⁵⁶⁷⁸⁹":
                if sup in text:
                    has_source_marker = True
                    break
            if "Quelle:" in text or "Source:" in text:
                has_source_marker = True
            # Bullets zählen (Zeilen mit •/·/-/* am Anfang)
            for line in text.split("\n"):
                ls = line.strip()
                if ls and len(ls) > 1 and ls[0] in ("•", "·") and ls[1] == " ":
                    bullets_count += 1
            # KPI-Numbers (Zahl + optional Einheit, oft als grosse Card-Text)
            if re.search(r'\b\d+([\.\']\d+)*\s*(%|Mio\.|k|CHF|EUR|h)?\b', text):
                has_kpi_numbers = True

        # Skip Title/Section/Closing-Folien für Bullet-Check
        is_structural = any(layout_name.startswith(p) for p in
                            ("Titelfolie", "Schlussfolie", "Abschlussfolie", "2_Zwischenfolie", "01_Agenda"))

        # 1. Generika-Check (unabhängig von KPI-Numbers — Footer enthält immer Folie-Nummer)
        slide_lower = slide_text.lower()
        for gw in _GENERIC_WORDS:
            if gw in slide_lower:
                findings.append({
                    "slide": slide_num,
                    "severity": "warn",
                    "issue": f"Generika-Verdacht: '{gw}' im Folientext",
                    "suggestion": f"'{gw}' mit konkreter Zahl/Name/Datum unterfüttern oder umformulieren",
                })
                break  # nur ein Finding pro Folie

        # 2. Bullet-Count auf Inhaltsfolien (nicht strukturell)
        if not is_structural and bullets_count > 0:
            if bullets_count <= 2:
                findings.append({
                    "slide": slide_num,
                    "severity": "warn",
                    "issue": f"Nur {bullets_count} Bullet(s) — Luftnummer-Verdacht",
                    "suggestion": "Min. 3 Bullets oder zusätzlicher Detail-/Kontext-Block",
                })
            elif bullets_count > 8:
                findings.append({
                    "slide": slide_num,
                    "severity": "warn",
                    "issue": f"{bullets_count} Bullets — Bleiwüsten-Verdacht",
                    "suggestion": "Auf 5-8 Bullets reduzieren oder auf 2 Folien splitten",
                })

        # 3. Quellen-Angabe auf Datenfolien
        if not is_structural and has_kpi_numbers and not has_source_marker:
            findings.append({
                "slide": slide_num,
                "severity": "info",
                "issue": "Daten ohne Quellen-Angabe",
                "suggestion": "footnote_block() mit Quelle ergänzen für MBB-Konformität",
            })

        # 4. Layout-Mix Tracking
        if layout_name == "04_vielText":
            add_content_count += 1

    # Globale Mix-Regel (Sprint 0 Density-Rules)
    inhaltsfolien_estimate = max(1, n_slides - 4)  # rough: minus Titel/Agenda/Section/Closing
    if add_content_count > 0:
        add_content_ratio = add_content_count / inhaltsfolien_estimate
        if add_content_ratio > 0.6:
            findings.append({
                "slide": 0,  # global
                "severity": "warn",
                "issue": f"{add_content_count}/{inhaltsfolien_estimate} Folien sind add_content ({add_content_ratio*100:.0f}%)",
                "suggestion": "Max. 60% add_content empfohlen — mehr visuelle Layouts (kpi_grid, two_column, flow_pipeline)",
            })

    # Output
    if findings and level != "silent":
        print(f"audit_deck: {len(findings)} findings")
        for f in findings:
            slide_str = f"global" if f['slide'] == 0 else f"Folie {f['slide']}"
            print(f"  {slide_str} [{f['severity']}]: {f['issue']}")
            print(f"    → {f['suggestion']}")
    if findings and level == "strict":
        raise ValueError(f"audit_deck(strict): {len(findings)} findings — siehe stdout für Details")
    return findings

def mece_check(bullet_list):
    """Heuristische MECE-Prüfung auf Bullet-Liste (Mutually Exclusive, Collectively Exhaustive).

    Args:
        bullet_list: list[str] — die zu prüfenden Bullets/Kategorien

    Returns:
        list[dict] mit Keys "bullet_idx", "bullet", "issue", "suggestion"

    Geprüfte Regeln:
        - Catch-All-Phrasen ("andere", "etc.", "u.a.") → nicht-MECE (collectively exhaustive verletzt)
        - Wort-Überlappung zwischen Bullets → mögliche Verletzung mutually-exclusive

    Beispiel:
        findings = mece_check([
            "Personal-Kosten",
            "IT-Personal-Aufwand",   # Überlappung mit "Personal"
            "Marketing",
            "Andere Kosten",         # Catch-All
        ])
        # → 2 Findings: Wort-Overlap + Catch-All
    """
    findings = []

    # 1. Catch-All-Phrasen
    for i, bullet in enumerate(bullet_list):
        bl = bullet.lower()
        for phrase in _CATCH_ALL_PHRASES:
            if phrase in bl:
                findings.append({
                    "bullet_idx": i,
                    "bullet": bullet,
                    "issue": f"Catch-All-Phrase '{phrase.strip()}'",
                    "suggestion": "Spezifische Items statt Sammel-Bullet, "
                                  "oder Sammel-Bullet explizit als 'Sonstige (X%)'",
                })
                break

    # 2. Wort-Überlappung
    def _significant_words(s):
        return set(w.lower().strip(".,;:!?") for w in s.split() if len(w) > 4)

    sigwords = [_significant_words(b) for b in bullet_list]
    for i in range(len(bullet_list)):
        for j in range(i + 1, len(bullet_list)):
            overlap = sigwords[i] & sigwords[j]
            if len(overlap) >= 2:
                overlap_str = ", ".join(sorted(overlap)[:3])
                findings.append({
                    "bullet_idx": i,
                    "bullet": bullet_list[i],
                    "issue": f"Wort-Überlappung mit Bullet {j + 1}: '{overlap_str}'",
                    "suggestion": "Kategorien klarer trennen — Mutually-Exclusive-Verstoss möglich",
                })
                break  # 1 Finding pro Bullet

    return findings

def action_title_check(headline):
    """Prüft ob Headline action-style (MBB-konform) oder Themen-Titel ist.

    Args:
        headline: str — die Headline (max 32 Zeichen empfohlen)

    Returns:
        dict mit Keys:
            "is_action":   bool — True wenn action-style
            "headline":    str — die geprüfte Headline
            "evidence":    list[str] — was die Klassifikation begründet
            "suggestion":  str — Verbesserungs-Vorschlag

    Heuristik (Action-Indikatoren):
        - Action-Verb erkannt (wächst, kippt, übersteigt, ...)
        - Zahl in der Headline (23%, CHF 240k, etc.)
        - Vergleichswort (vs, statt, über, unter)

    Beispiele:
        action_title_check("Print noch dominant")
        # → is_action=True (Verb "ist" implizit + "dominant" als Aussage)
        action_title_check("EBIT-Entwicklung")
        # → is_action=False (Themen-Titel, kein Verb, keine Zahl)
        action_title_check("EBIT wächst 23% YoY")
        # → is_action=True (Verb + Zahl)
    """
    h = headline.strip()
    hl = h.lower()
    evidence = []

    # Action-Verb
    has_verb = any(v in hl.split() or v in hl for v in _ACTION_VERBS)
    if has_verb:
        evidence.append("Action-Verb erkannt")

    # ECHTE Aussage-Zahl: number + Unit (NICHT nur "Q1" oder Jahreszahl)
    # Beispiele die zählen: "23%", "CHF 240k", "18 Mt", "3'900h"
    # Beispiele die NICHT zählen: "Q1", "2026", "Ergebnis 2025"
    has_real_number = bool(re.search(
        r'\d+\s*(%|‰|Mio\.|Mrd\.|CHF|EUR|USD|\$|€|k\b|h\b|MW|kW|FTE|Pkt\.)',
        h, re.IGNORECASE))
    if has_real_number:
        evidence.append("Zahl mit Einheit")

    # Prozent (separate Prüfung für saubere Evidence)
    if "%" in h and "Zahl mit Einheit" not in evidence:
        evidence.append("Prozent-Wert")

    # Währungs-Kontext (nur wenn nicht schon durch has_real_number erfasst)
    if any(c in h for c in ["CHF", "€", "$", "EUR", "USD"]) and "Zahl mit Einheit" not in evidence:
        evidence.append("Währungs-Wert")

    # Vergleichs-Operator
    has_comp = any(p in hl for p in [" vs ", " vs.", " statt ", " gegen ", " über ", " unter ", " mehr als ", " weniger als "])
    if has_comp:
        evidence.append("Vergleich/Kontrast")

    # Adjektiv-Aussage (z.B. "stark", "schwach", "dominant", "knapp")
    AUSSAGE_ADJ = {"dominant", "stark", "schwach", "knapp", "üppig", "kritisch",
                    "rentabel", "unrentabel", "profitabel", "verlustreich",
                    "wichtig", "entscheidend", "tragfähig", "nachhaltig"}
    if any(adj in hl.split() for adj in AUSSAGE_ADJ):
        evidence.append("Aussage-Adjektiv")

    is_action = bool(evidence)

    if is_action:
        suggestion = "Action-Title — gut. Falls noch klarer: konkrete Zahl ergänzen."
    else:
        suggestion = (
            "Themen-Titel erkannt. Empfehlung: konkrete Aussage statt Thema. "
            "Beispiel: 'EBIT wächst 23%' statt 'Ergebnis Q1', "
            "'Print dreht Q3' statt 'Print-Entwicklung'."
        )

    return {
        "is_action": is_action,
        "headline": h,
        "evidence": evidence,
        "suggestion": suggestion,
    }

# ── Sprint 12: Density Discipline & Sub-Section-Tracker ──────────────────────

def add_subsection_tracker(slide, sections, current_section, current_subsection):
    """Mini-Roadmap am Folienkopf: Hauptsektion + Subsektion-Fortschritt.
    Sprint 12 — siehe ROADMAP.md"""
    _stub("add_subsection_tracker", 12)

def add_density_meter(slide, content_stats):
    """Author-Tool: visualisiert Folien-Dichte. Im Render unsichtbar.
    Sprint 12 — siehe ROADMAP.md"""
    _stub("add_density_meter", 12)

def enforce_density(slide, min_elements=4, max_elements=12):
    """Raise wenn Folie unter Mindest- oder ueber Maximum-Inhaltsdichte.
    Sprint 12 — siehe ROADMAP.md"""
    _stub("enforce_density", 12)

# ── Sprint 13: Action-Titles (optional) ──────────────────────────────────────

def add_content_action_title(prs, kicker, action_title, body, folio="", source=""):
    """Inhaltsfolie mit 60-80-Zeichen Action-Title (statt 32-Zeichen-Limit).
    Adaptive Font-Size: 1-Zeile 72pt bis 32 Zeichen, danach 2-Zeilen 54pt.
    NUR fuer Inhaltsfolien — Cover/Section/Closing bleiben bei 32 Zeichen.
    Sprint 13 (optional) — siehe ROADMAP.md"""
    _stub("add_content_action_title", 13)
